"""Rich document creation — markdown-aware writers + the convert_document tool.

String content now parses as structured markdown (documents/markdown.py) and
renders as REAL headings/lists/tables in .docx/.pdf/.pptx/.html, while flat
plain-text strings must keep writing fine everywhere (backward compat).
"""

from __future__ import annotations

import csv
from pathlib import Path

from iron_jarvis.documents import (
    SUPPORTED_WRITE,
    document_tools,
    extract_text,
    parse_markdown,
    write_document,
)
from iron_jarvis.tools.base import ToolContext

MD = """# Annual Report

Intro paragraph with **bold** and *italic* words.

## Numbers

- first bullet
  - nested bullet
1. step one
2. step two

| name | qty |
| ---- | --- |
| arc reactor | 3 |
| suit | 42 |

```
code_line = 1
```
"""


def _ctx(workspace: Path) -> ToolContext:
    """Minimal ToolContext — only ``.workspace`` is used by the document tools."""
    return ToolContext(
        workspace=workspace,
        session_id="t",
        agent_run_id="t",
        config=None,
        event_bus=None,
        engine=None,
    )


def _tool(name: str):
    return next(t for t in document_tools() if t.name == name)


# --- the block model ------------------------------------------------------------


def test_parse_markdown_blocks():
    blocks = parse_markdown(MD)
    kinds = [b.kind for b in blocks]
    assert kinds == [
        "heading",
        "paragraph",
        "heading",
        "bullet",
        "bullet",
        "numbered",
        "numbered",
        "table",
        "code",
    ]
    assert blocks[0].level == 1 and blocks[0].text == "Annual Report"
    assert ("bold", True, False) in blocks[1].runs
    assert ("italic", False, True) in blocks[1].runs
    assert blocks[4].level == 1  # nested bullet
    assert blocks[7].rows == [["name", "qty"], ["arc reactor", "3"], ["suit", "42"]]
    assert blocks[8].text == "code_line = 1"


def test_parse_markdown_never_raises():
    for weird in ("", None, "| |", "```unclosed\nx", "***", "  \n\t\n", 42):
        assert isinstance(parse_markdown(weird), list)


# --- docx -----------------------------------------------------------------------


def test_md_to_docx_rich(tmp_path):
    import docx

    p = tmp_path / "rich.docx"
    write_document(p, MD)
    d = docx.Document(str(p))

    styles = [para.style.name for para in d.paragraphs]
    assert "Heading 1" in styles
    assert "Heading 2" in styles
    assert "List Bullet" in styles
    assert "List Bullet 2" in styles  # nested bullet
    assert "List Number" in styles

    # A REAL table, with a bold header row.
    assert len(d.tables) == 1
    t = d.tables[0]
    assert len(t.rows) == 3
    assert t.cell(0, 0).text == "name"
    assert t.cell(1, 0).text == "arc reactor"
    assert t.cell(0, 0).paragraphs[0].runs[0].bold is True

    # Inline bold run in the intro paragraph, and monospace code.
    all_runs = [r for para in d.paragraphs for r in para.runs]
    assert any(r.text == "bold" and r.bold for r in all_runs)
    assert any(r.font.name == "Consolas" for r in all_runs)


# --- pdf ------------------------------------------------------------------------


def test_md_to_pdf_parseable(tmp_path):
    from pypdf import PdfReader

    p = tmp_path / "rich.pdf"
    write_document(p, MD)
    text = "\n".join((pg.extract_text() or "") for pg in PdfReader(str(p)).pages)
    assert "Annual Report" in text
    assert "first bullet" in text
    assert "arc reactor" in text


# --- pptx -----------------------------------------------------------------------


def test_md_to_pptx_sections(tmp_path):
    import pptx

    p = tmp_path / "deck.pptx"
    write_document(p, "# One\n\n- a\n  - a nested\n- b\n\n# Two\n\npoint")
    prs = pptx.Presentation(str(p))
    # 2 '#' sections -> title slide + one slide per section == 3.
    assert len(prs.slides) == 3
    assert prs.slides[1].shapes.title.text == "One"
    assert prs.slides[2].shapes.title.text == "Two"
    body = prs.slides[1].placeholders[1].text_frame
    assert [para.text for para in body.paragraphs] == ["a", "a nested", "b"]
    assert body.paragraphs[1].level == 1  # nested bullet indented one level


def test_pptx_hr_starts_slide(tmp_path):
    import pptx

    p = tmp_path / "hr.pptx"
    write_document(p, "# Intro\n\nhello\n\n---\n\nsecond section line")
    prs = pptx.Presentation(str(p))
    assert len(prs.slides) == 3


def test_pptx_without_sections_keeps_legacy_deck(tmp_path):
    import pptx

    p = tmp_path / "flat.pptx"
    write_document(p, "Quarterly Review\nRevenue up\nCosts down")
    prs = pptx.Presentation(str(p))
    assert len(prs.slides) == 2  # today's behavior: title + one bullet slide
    assert prs.slides[0].shapes.title.text == "Quarterly Review"


# --- xlsx -----------------------------------------------------------------------


def test_xlsx_multi_sheet_dict(tmp_path):
    from openpyxl import load_workbook

    p = tmp_path / "book.xlsx"
    write_document(
        p,
        {
            "sheets": {
                "Data": [["name", "qty"], ["arc reactor", 3], ["suit", 42]],
                "Calc": [["label", "value"], ["total", "=SUM(1,2)"]],
            }
        },
    )
    wb = load_workbook(str(p))  # data_only=False so formulas stay visible
    assert wb.sheetnames == ["Data", "Calc"]

    data = wb["Data"]
    assert data["A1"].value == "name"
    assert data["A1"].font.bold is True  # header row bolded
    assert data.freeze_panes == "A2"
    assert 8 <= data.column_dimensions["A"].width <= 60

    assert wb["Calc"]["B2"].value == "=SUM(1,2)"  # formula NOT escaped


def test_xlsx_list_rows_unchanged(tmp_path):
    p = tmp_path / "rows.xlsx"
    write_document(p, [["a", "b"], ["1", "2"]])
    text = extract_text(p)
    for token in ("a", "b", "1", "2"):
        assert token in text


def test_xlsx_string_content_unchanged(tmp_path):
    p = tmp_path / "lines.xlsx"
    write_document(p, "# not markdown here\n- still a raw line")
    text = extract_text(p)
    assert "# not markdown here" in text
    assert "- still a raw line" in text


# --- html -----------------------------------------------------------------------


def test_html_rich(tmp_path):
    assert ".html" in SUPPORTED_WRITE and ".htm" in SUPPORTED_WRITE
    p = tmp_path / "page.html"
    write_document(p, MD)
    html = p.read_text(encoding="utf-8")
    assert "<h1>Annual Report</h1>" in html
    assert "<table" in html and "<th>name</th>" in html
    assert "<strong>bold</strong>" in html
    assert "<em>italic</em>" in html
    assert "<li>first bullet</li>" in html
    assert "<pre><code>" in html


def test_html_passthrough_for_full_pages(tmp_path):
    p = tmp_path / "raw.html"
    page = "<!DOCTYPE html><html><body><p>already html</p></body></html>"
    write_document(p, page)
    assert p.read_text(encoding="utf-8") == page


# --- convert_document tool --------------------------------------------------------


async def test_convert_csv_xlsx_csv_roundtrip(tmp_path):
    ws = tmp_path / "ws"
    ws.mkdir()
    ctx = _ctx(ws)
    rows = [
        ["name", "note"],
        ["Tony", "likes, commas"],
        ["Pepper", 'quotes "inside"'],
    ]
    with open(ws / "src.csv", "w", newline="", encoding="utf-8") as f:
        csv.writer(f).writerows(rows)

    tool = _tool("convert_document")
    res1 = await tool.execute({"source": "src.csv", "target": "mid.xlsx"}, ctx)
    assert res1.ok, res1.error
    res2 = await tool.execute({"source": "mid.xlsx", "target": "back.csv"}, ctx)
    assert res2.ok, res2.error

    with open(ws / "back.csv", newline="", encoding="utf-8") as f:
        assert [list(r) for r in csv.reader(f)] == rows


async def test_convert_docx_to_pdf(tmp_path):
    ws = tmp_path / "ws"
    ws.mkdir()
    ctx = _ctx(ws)
    write_document(ws / "memo.docx", "conversion fidelity marker")

    res = await _tool("convert_document").execute(
        {"source": "memo.docx", "target": "memo.pdf"}, ctx
    )
    assert res.ok, res.error
    assert "marker" in extract_text(ws / "memo.pdf")


async def test_convert_rejects_unsupported_suffixes(tmp_path):
    ws = tmp_path / "ws"
    ws.mkdir()
    ctx = _ctx(ws)
    (ws / "note.txt").write_text("hi", encoding="utf-8")

    bad_target = await _tool("convert_document").execute(
        {"source": "note.txt", "target": "out.exe"}, ctx
    )
    assert not bad_target.ok
    assert "supported target formats" in bad_target.error

    bad_source = await _tool("convert_document").execute(
        {"source": "blob.bin", "target": "out.txt"}, ctx
    )
    assert not bad_source.ok
    assert "supported source formats" in bad_source.error


async def test_convert_missing_source_is_graceful(tmp_path):
    ws = tmp_path / "ws"
    ws.mkdir()
    res = await _tool("convert_document").execute(
        {"source": "nope.csv", "target": "out.xlsx"}, _ctx(ws)
    )
    assert not res.ok
    assert res.error


# --- backward compat: flat plain text writes fine everywhere ---------------------


def test_plain_text_compat_every_format(tmp_path):
    for suffix in (".docx", ".xlsx", ".pptx", ".pdf", ".csv", ".txt", ".md", ".html"):
        p = tmp_path / f"plain{suffix}"
        write_document(p, "just some plain text\nsecond line")  # must not raise
        assert p.is_file() and p.stat().st_size > 0
        assert "plain" in extract_text(p)


def test_pdf_full_unicode_roundtrip(tmp_path):
    """Bundled DejaVu fonts => PDFs keep real unicode (no more latin-1 mangling)."""
    from pypdf import PdfReader

    from iron_jarvis.documents.writers import write_document

    md = "# Résumé — München\n\nCafé naïve • Ω ≈ 3.14 → done\n\n- Ærøskøbing\n"
    p = write_document(tmp_path / "u.pdf", md)
    text = "".join(page.extract_text() or "" for page in PdfReader(str(p)).pages)
    for needle in ("Résumé", "München", "naïve", "Ærøskøbing"):
        assert needle in text, f"missing {needle!r} in extracted PDF text"


def test_pdf_falls_back_when_fonts_missing(tmp_path, monkeypatch):
    """A stripped install (no TTFs) must still write valid PDFs via Latin-1."""
    import iron_jarvis.documents.writers as w
    from pypdf import PdfReader

    monkeypatch.setattr(w, "_FONT_DIR", tmp_path / "nope")
    p = w.write_document(tmp_path / "f.pdf", "# Hello\n\nplain fallback")
    text = "".join(page.extract_text() or "" for page in PdfReader(str(p)).pages)
    assert "Hello" in text
