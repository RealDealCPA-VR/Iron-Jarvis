"""Reliability hardening for office file READING (documents/readers.py).

Covers the daily-driver failure modes that used to corrupt or silently swallow
real content: legacy-codepage decoding, scanned/empty PDFs, encrypted files,
xlsx date/formula rendering, docx headers/footers, pptx notes, page/sheet
slicing, RTF and named legacy formats. All offline.

Non-ASCII test strings are built from \\u escapes so the source file stays pure
ASCII (some transfer paths mangle raw non-ASCII bytes in .py files).
"""

from __future__ import annotations

import datetime

import pytest

from iron_jarvis.documents.readers import (
    _SCANNED_PDF_SENTINEL,
    _decode_bytes,
    extract_text,
)

# "Cafe -- naive resume" — cp1252 accents (0xE9/0xEF) + em dash (0x97), built
# from escapes so the .py file is pure ASCII and cannot be mangled in transit.
_ACCENTED = "Café — naïve résumé"
_REPLACEMENT = "�"  # U+FFFD, what a botched utf-8 decode would produce


# --- encoding detection -------------------------------------------------------


def test_cp1252_roundtrip(tmp_path):
    p = tmp_path / "legacy.txt"
    p.write_bytes(_ACCENTED.encode("cp1252"))
    assert extract_text(p) == _ACCENTED
    assert _REPLACEMENT not in extract_text(p)


def test_utf8_bom_stripped(tmp_path):
    # Excel/Windows CSV exports commonly prepend a UTF-8 BOM.
    p = tmp_path / "bom.txt"
    p.write_bytes(b"\xef\xbb\xbfhello")
    assert extract_text(p) == "hello"


def test_csv_cp1252(tmp_path):
    p = tmp_path / "data.csv"
    p.write_bytes(("name,note\nTony," + _ACCENTED + "\n").encode("cp1252"))
    out = extract_text(p)
    assert _ACCENTED in out
    assert _REPLACEMENT not in out


def test_decode_bytes_latin1_safety_net():
    # A byte undefined in cp1252 (0x81) must still decode, not raise.
    assert _decode_bytes(b"a\x81b")  # latin-1 maps every byte


# --- scanned / empty sentinel -------------------------------------------------


def test_scanned_pdf_sentinel(tmp_path):
    # A valid PDF with no extractable text (no text operators) => sentinel, not
    # an ambiguous empty string, so an agent knows it is unreadable.
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    p = tmp_path / "scanned.pdf"
    with open(p, "wb") as f:
        writer.write(f)
    assert extract_text(p) == _SCANNED_PDF_SENTINEL


# --- encrypted files ----------------------------------------------------------


def test_encrypted_pdf_clear_error(tmp_path):
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.encrypt("secret")  # non-empty owner+user password
    p = tmp_path / "locked.pdf"
    with open(p, "wb") as f:
        writer.write(f)
    with pytest.raises(ValueError, match="password-protected"):
        extract_text(p)


def test_encrypted_office_ole_magic(tmp_path):
    # A password-protected .xlsx is an OLE/CFB container, not a zip. We must give
    # a clear error instead of the misleading "BadZipFile: File is not a zip".
    p = tmp_path / "locked.xlsx"
    p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 32)
    with pytest.raises(ValueError, match="password-protected|encrypted"):
        extract_text(p)


# --- xlsx dates + formulas ----------------------------------------------------


def _make_xlsx(path, build):
    from openpyxl import Workbook

    wb = Workbook()
    build(wb)
    wb.save(str(path))


def test_xlsx_date_formatting(tmp_path):
    def build(wb):
        ws = wb.active
        d = ws["A1"]
        d.value = datetime.datetime(2024, 1, 2)  # midnight
        d.number_format = "yyyy-mm-dd"
        ws["B1"].value = datetime.datetime(2024, 1, 2, 13, 30)  # has a time

    p = tmp_path / "dates.xlsx"
    _make_xlsx(p, build)
    out = extract_text(p)
    assert "2024-01-02" in out
    assert "00:00:00" not in out  # date-only must NOT gain a spurious time
    assert "13:30" in out  # a real datetime keeps its time


def test_xlsx_formula_fallback(tmp_path):
    # openpyxl-written formulas have no cached value, so data_only yields None.
    # We must fall back to showing the formula text, not a blank cell.
    def build(wb):
        ws = wb.active
        ws["A1"].value = 2
        ws["B1"].value = 3
        ws["C1"].value = "=A1+B1"

    p = tmp_path / "calc.xlsx"
    _make_xlsx(p, build)
    out = extract_text(p)
    assert "=A1+B1" in out


def test_xlsx_sheet_selection(tmp_path):
    def build(wb):
        wb.active.title = "First"
        wb.active["A1"].value = "alpha"
        two = wb.create_sheet("Second")
        two["A1"].value = "beta"

    p = tmp_path / "multi.xlsx"
    _make_xlsx(p, build)
    by_name = extract_text(p, sheet="Second")
    assert "beta" in by_name and "alpha" not in by_name
    by_index = extract_text(p, sheet=1)
    assert "beta" in by_index
    with pytest.raises(ValueError, match="not found"):
        extract_text(p, sheet="Nope")


# --- docx headers / footers / tables ------------------------------------------


def test_docx_headers_footers_and_inline_table(tmp_path):
    import docx

    d = docx.Document()
    d.add_paragraph("Body prose")
    d.sections[0].header.paragraphs[0].text = "LETTERHEAD"
    d.sections[0].footer.paragraphs[0].text = "PAGE FOOTER"
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "cellA"
    t.rows[0].cells[1].text = "cellB"
    p = tmp_path / "letter.docx"
    d.save(str(p))
    out = extract_text(p)
    for token in ("LETTERHEAD", "Body prose", "cellA", "cellB", "PAGE FOOTER"):
        assert token in out


# --- pptx notes + tables ------------------------------------------------------


def test_pptx_notes_and_table(tmp_path):
    import pptx
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "DeckTitle"
    slide.notes_slide.notes_text_frame.text = "SpeakerNote"
    tbl = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
    ).table
    tbl.cell(0, 0).text = "H1"
    tbl.cell(1, 1).text = "V22"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))
    out = extract_text(p)
    for token in ("DeckTitle", "SpeakerNote", "H1", "V22"):
        assert token in out


# --- page-range slicing -------------------------------------------------------


def test_pdf_page_range(tmp_path):
    from iron_jarvis.documents import write_document

    p = tmp_path / "multi.pdf"
    write_document(p, "PageOneMarker")
    assert "PageOneMarker" in extract_text(p, page_range="1")
    # Out-of-range slice selects no pages => scanned sentinel (nothing extracted).
    assert extract_text(p, page_range="99") == _SCANNED_PDF_SENTINEL


# --- rtf + legacy -------------------------------------------------------------


def test_rtf_reader_strips_control_words(tmp_path):
    p = tmp_path / "note.rtf"
    p.write_text(r"{\rtf1\ansi\deff0 Hello \b world\b0 .\par}", encoding="utf-8")
    out = extract_text(p)
    assert "Hello" in out and "world" in out
    assert "\\rtf1" not in out and "\\par" not in out


@pytest.mark.parametrize(
    "suffix,needle",
    [(".doc", "docx"), (".xls", "xlsx"), (".ppt", "pptx"), (".odt", "docx")],
)
def test_legacy_formats_clear_error(tmp_path, suffix, needle):
    p = tmp_path / f"old{suffix}"
    p.write_bytes(b"\xd0\xcf\x11\xe0legacy binary")
    with pytest.raises(ValueError, match=needle):
        extract_text(p)


# --- size guard ---------------------------------------------------------------


def test_size_guard(tmp_path, monkeypatch):
    import iron_jarvis.documents.readers as readers

    monkeypatch.setattr(readers, "_MAX_BYTES", 8)
    p = tmp_path / "big.txt"
    p.write_text("this is more than eight bytes", encoding="utf-8")
    with pytest.raises(ValueError, match="too large"):
        extract_text(p)
