"""High-fidelity document creation — the reliability improvements.

Covers the new writer/parser behaviors: inline links/images/code, valid JSON
serialization of dict content, pptx overflow pagination, paragraph joining of
soft-wrapped lines, atomic writes that preserve the old file on failure, and
per-column table alignment captured from the ``|:--:|`` separator row.
"""

from __future__ import annotations

import json

import pytest

from iron_jarvis.documents import parse_markdown, write_document
from iron_jarvis.documents import writers as w


# --- inline links / images / code (no literal brackets ever) -------------------


def test_inline_link_parsing_leaves_no_brackets():
    blocks = parse_markdown("See [the Docs](https://example.com/x) today.")
    runs = blocks[0].runs
    # No run carries the raw markdown markup.
    assert all("](" not in r[0] and "[" not in r[0] for r in runs)
    link = next(r for r in runs if getattr(r, "href", None))
    assert link.href == "https://example.com/x"
    assert link[0] == "the Docs"
    # And the flattened plain text is clean prose.
    assert "[" not in blocks[0].text and "](" not in blocks[0].text


def test_inline_image_and_code_parsed():
    runs = parse_markdown("![logo](http://a/b.png) and `x=1`")[0].runs
    img = next(r for r in runs if getattr(r, "image", False))
    assert img.href == "http://a/b.png" and img[0] == "logo"
    code = next(r for r in runs if getattr(r, "code", False))
    assert code[0] == "x=1"


def test_link_renders_in_html(tmp_path):
    p = tmp_path / "l.html"
    write_document(p, "Read [our guide](https://example.com/g).")
    html = p.read_text(encoding="utf-8")
    assert '<a href="https://example.com/g">our guide</a>' in html
    assert "](" not in html  # never a literal markdown link


def test_link_renders_in_docx_without_brackets(tmp_path):
    import docx

    p = tmp_path / "l.docx"
    write_document(p, "Visit [site](https://example.com/site) now.")
    d = docx.Document(str(p))
    # The hyperlink target is stored as an external relationship.
    targets = [r.target_ref for r in d.part.rels.values()]
    assert "https://example.com/site" in targets
    # The document XML must not contain the literal ``](`` markup.
    assert "](" not in d.element.xml


def test_link_renders_in_pdf(tmp_path):
    from pypdf import PdfReader

    p = tmp_path / "l.pdf"
    write_document(p, "Open [the report](https://example.com/r) please.")
    text = "".join(pg.extract_text() or "" for pg in PdfReader(str(p)).pages)
    assert "the report" in text
    assert "](" not in text


# --- json / yaml serialization -------------------------------------------------


def test_json_dict_serializes_valid(tmp_path):
    p = tmp_path / "data.json"
    obj = {"name": "Tony", "vals": [1, 2, 3], "flag": True, "nil": None}
    write_document(p, obj)
    raw = p.read_text(encoding="utf-8")
    # Must be REAL json, not a Python repr (which the old text path produced).
    assert "'" not in raw  # no single-quoted keys
    assert "true" in raw and "null" in raw  # json literals, not True/None
    assert json.loads(raw) == obj


def test_json_string_passthrough(tmp_path):
    p = tmp_path / "s.json"
    write_document(p, '{"already": "serialized"}')
    assert json.loads(p.read_text(encoding="utf-8")) == {"already": "serialized"}


def test_yaml_dict_serializes(tmp_path):
    p = tmp_path / "c.yaml"
    write_document(p, {"a": 1, "b": ["x", "y"]})
    raw = p.read_text(encoding="utf-8")
    # Whether via PyYAML or the JSON fallback, it must round-trip to the dict.
    try:
        import yaml

        assert yaml.safe_load(raw) == {"a": 1, "b": ["x", "y"]}
    except ImportError:
        assert json.loads(raw) == {"a": 1, "b": ["x", "y"]}


# --- pptx overflow pagination --------------------------------------------------


def test_pptx_overflow_paginates(tmp_path):
    import pptx

    bullets = "\n".join(f"- item {i}" for i in range(20))
    p = tmp_path / "big.pptx"
    write_document(p, f"# Big Section\n\n{bullets}")
    prs = pptx.Presentation(str(p))
    titles = [s.shapes.title.text for s in prs.slides if s.shapes.title]
    # 20 items / ~9 per slide => 3 content slides + the title slide.
    assert len(prs.slides) >= 4
    # Overflow spilled onto at least one continuation slide.
    assert any(t == "Big Section (cont.)" for t in titles)
    # No single body frame silently holds all 20 bullets.
    max_bullets = max(
        len(s.placeholders[1].text_frame.paragraphs)
        for s in list(prs.slides)[1:]
        if len(s.placeholders) > 1
    )
    assert max_bullets <= 9


# --- paragraph joining ---------------------------------------------------------


def test_paragraph_joins_soft_wrapped_lines():
    blocks = parse_markdown("line one\nline two\nline three\n\nsecond para")
    paras = [b for b in blocks if b.kind == "paragraph"]
    assert len(paras) == 2
    assert paras[0].text == "line one line two line three"
    assert paras[1].text == "second para"


def test_paragraph_joining_stops_at_structure():
    blocks = parse_markdown("intro text\n- a bullet\nmore text")
    kinds = [b.kind for b in blocks]
    assert kinds == ["paragraph", "bullet", "paragraph"]
    assert blocks[0].text == "intro text"


# --- atomic write --------------------------------------------------------------


def test_atomic_write_preserves_old_file_on_failure(tmp_path):
    p = tmp_path / "keep.txt"
    p.write_text("ORIGINAL", encoding="utf-8")

    with pytest.raises(RuntimeError):
        with w._atomic(p) as tmp:
            tmp.write_text("HALF-WRITTEN", encoding="utf-8")
            raise RuntimeError("mid-save failure")

    # The good file is untouched and no temp turd is left behind.
    assert p.read_text(encoding="utf-8") == "ORIGINAL"
    assert list(tmp_path.iterdir()) == [p]


def test_atomic_write_replaces_on_success(tmp_path):
    p = tmp_path / "out.txt"
    p.write_text("OLD", encoding="utf-8")
    write_document(p, "NEW CONTENT")
    assert p.read_text(encoding="utf-8") == "NEW CONTENT"
    assert list(tmp_path.iterdir()) == [p]


# --- table alignment -----------------------------------------------------------


def test_table_alignment_captured():
    md = "| a | b | c |\n|:--|:-:|--:|\n| 1 | 2 | 3 |\n"
    table = next(b for b in parse_markdown(md) if b.kind == "table")
    assert table.aligns == ["left", "center", "right"]
    assert table.rows == [["a", "b", "c"], ["1", "2", "3"]]


def test_table_alignment_in_html(tmp_path):
    p = tmp_path / "t.html"
    write_document(p, "| a | b |\n|:-:|--:|\n| 1 | 2 |\n")
    html = p.read_text(encoding="utf-8")
    assert "text-align:center" in html
    assert "text-align:right" in html


def test_gfm_table_without_outer_pipes():
    md = "Col A | Col B\n--- | ---\nx | y\n"
    table = next(b for b in parse_markdown(md) if b.kind == "table")
    assert table.rows == [["Col A", "Col B"], ["x", "y"]]


# --- xlsx numeric / date coercion ----------------------------------------------


def test_xlsx_data_cell_coercion(tmp_path):
    from openpyxl import load_workbook

    p = tmp_path / "c.xlsx"
    write_document(
        p,
        {
            "sheets": {
                "S": [
                    ["id", "amount", "when"],
                    ["007", "42", "2026-07-10"],
                    ["012", "3.5", "=SUM(1,2)"],
                ]
            }
        },
    )
    ws = load_workbook(str(p))["S"]
    assert ws["A1"].value == "id"  # header stays text
    assert ws["A2"].value == "007"  # leading-zero id stays text
    assert ws["B2"].value == 42 and isinstance(ws["B2"].value, int)
    assert ws["B3"].value == 3.5 and isinstance(ws["B3"].value, float)
    assert str(ws["C2"].value).startswith("2026-07-10")  # ISO date -> datetime
    assert ws["C3"].value == "=SUM(1,2)"  # formula preserved
