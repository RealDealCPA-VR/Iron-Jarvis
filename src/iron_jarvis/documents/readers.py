"""Document text extraction (readers).

``extract_text(path)`` reads a real-world file and returns its text content,
dispatching on the lowercased filename suffix:

* ``.pdf``  -> pypdf (join every page's ``extract_text()``); optional ``page_range``.
* ``.docx`` -> python-docx (headers/footers + body in document order + tables +
  footnotes/endnotes).
* ``.xlsx`` -> openpyxl (typed date formatting + formula fallback); optional ``sheet``.
* ``.pptx`` -> python-pptx (shapes + tables + speaker notes + groups); optional
  ``page_range`` (slides).
* ``.csv``  -> stdlib csv (tab-joined rows), encoding-detected.
* ``.rtf``  -> striprtf if available, else a control-word stripper.
* ``.txt/.md/.json/.py/.js/.ts/.html/.yaml/.yml/.log`` and any unknown-but-text
  file -> encoding-detected decode (utf-8-sig → charset-normalizer → cp1252 →
  latin-1, so cp1252/latin-1 office exports survive instead of turning into
  replacement chars).
* ``.png/.jpg/.jpeg/.gif/.bmp/.webp`` -> Pillow, returning a concise note such as
  ``"[image PNG 800x600, mode RGB]"`` (NO OCR).

Reliability guarantees the office-daily-driver path depends on:

* Encrypted PDFs and OLE/CFB-wrapped (password-protected) Office files raise a
  CLEAR error instead of a cryptic ``BadZipFile`` / garbage.
* A PDF/docx/pptx that extracts to effectively nothing returns a sentinel string
  so an agent can distinguish "scanned/image-only, unreadable" from "genuinely
  empty" rather than an ambiguous ``""``.
* Files above ~100 MB short-circuit with a clear error (PDF/text/unknown paths)
  instead of blowing up memory.

A clear :class:`ValueError` is raised for a genuinely unsupported / binary type
or a named legacy format (.doc/.xls/.ppt/.odt) that needs conversion first.
"""

from __future__ import annotations

import csv
import datetime as _dt
import io
import re
from collections.abc import Iterable
from pathlib import Path

#: A stray leading BOM (U+FEFF) survives some decodes; strip it so callers never
#: see it as the first character of otherwise-clean text.
_FEFF = "﻿"

#: Hard ceiling for the whole-file read paths (PDF/text/unknown). Reading the
#: user's real files must never OOM the daemon; slices go through page_range.
_MAX_BYTES = 100 * 1024 * 1024

#: Returned when a PDF yields no extractable text — almost always a scanned /
#: image-only document. Kept as an exact, greppable sentinel so agents (and the
#: dashboard) can branch on "unreadable" vs. "genuinely empty".
_SCANNED_PDF_SENTINEL = (
    "[no extractable text — likely a scanned/image-only PDF; OCR not available]"
)

#: OLE2 / Compound File Binary magic. Encrypted OOXML (a password-protected
#: .docx/.xlsx/.pptx) is wrapped in this container, so python-docx/openpyxl see
#: it as "not a zip file". Detect it up front to give an honest error.
_OLE_MAGIC = b"\xd0\xcf\x11\xe0"

#: Suffixes read verbatim as encoding-detected text.
_TEXT_SUFFIXES: frozenset[str] = frozenset(
    {
        ".txt",
        ".md",
        ".json",
        ".py",
        ".js",
        ".ts",
        ".html",
        ".htm",
        ".yaml",
        ".yml",
        ".log",
    }
)

#: Raster image suffixes -> described, never OCR'd.
_IMAGE_SUFFIXES: frozenset[str] = frozenset(
    {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
)

#: Structured document suffixes handled by dedicated parsers.
_DOC_SUFFIXES: frozenset[str] = frozenset(
    {".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".rtf"}
)

#: Legacy binary Office / ODF formats we cannot parse — mapped to an honest
#: error naming the format and the modern target to convert to.
_LEGACY_HINTS: dict[str, str] = {
    ".doc": "legacy Word 97-2003 (.doc) — convert it to .docx first",
    ".xls": "legacy Excel 97-2003 (.xls) — convert it to .xlsx first",
    ".ppt": "legacy PowerPoint 97-2003 (.ppt) — convert it to .pptx first",
    ".odt": "OpenDocument Text (.odt) — convert it to .docx first",
    ".ods": "OpenDocument Spreadsheet (.ods) — convert it to .xlsx first",
    ".odp": "OpenDocument Presentation (.odp) — convert it to .pptx first",
}

#: Every suffix ``extract_text`` knows how to read (unknown text files also work,
#: but are not advertised here).
SUPPORTED_READ: set[str] = set(_DOC_SUFFIXES | _TEXT_SUFFIXES | _IMAGE_SUFFIXES)


def extract_text(
    path: str | Path,
    *,
    page_range: str | None = None,
    sheet: str | int | None = None,
) -> str:
    """Return the text content of ``path``, dispatched by file suffix.

    ``page_range`` (e.g. ``"2"``, ``"1-3"``, ``"2-"``, ``"1,4-6"``; 1-based,
    inclusive) slices PDF pages / PPTX slides. ``sheet`` (name or 0-based index)
    selects a single worksheet of an XLSX. Both are ignored by formats they do
    not apply to.

    Raises :class:`ValueError` for a truly unsupported / binary file type (or a
    password-protected / oversized / legacy file) and :class:`FileNotFoundError`
    if the path does not exist.
    """
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"no such file: {p}")
    if p.is_dir():
        raise ValueError(f"path is a directory, not a document: {p}")

    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(p, page_range=page_range)
    if suffix == ".docx":
        return _read_docx(p)
    if suffix == ".xlsx":
        return _read_xlsx(p, sheet=sheet)
    if suffix == ".pptx":
        return _read_pptx(p, page_range=page_range)
    if suffix == ".csv":
        return _read_csv(p)
    if suffix == ".rtf":
        return _read_rtf(p)
    if suffix in _IMAGE_SUFFIXES:
        return _describe_image(p)
    if suffix in _LEGACY_HINTS:
        # Named legacy format: fail with a useful next step, not "binary file".
        raise ValueError(f"cannot read {_LEGACY_HINTS[suffix]}")
    if suffix in _TEXT_SUFFIXES:
        _check_size(p, "text file")
        return _decode_bytes(p.read_bytes())

    # Unknown suffix: sniff for binary content; decode if it looks like text.
    return _read_unknown(p)


# --- encoding + size helpers --------------------------------------------------


def _decode_bytes(data: bytes) -> str:
    """Best-effort decode of arbitrary text bytes, never raising.

    Order matters and is deliberate: utf-8 (with BOM stripping) is correct for
    the vast majority and for Excel/Windows CSV exports that prepend a UTF-8 BOM;
    cp1252 (strict) is the single most common Windows legacy codepage and MUST
    win over statistical guessing — charset-normalizer notoriously mis-detects
    short cp1252 strings as Latin-2 (turning 0xEF 'ï' into 'ď'), which is exactly
    the corruption we are here to prevent; charset-normalizer therefore runs only
    for bytes cp1252 rejects (its 5 undefined code points), where a smart guess
    genuinely helps; latin-1 maps all 256 byte values so it never raises and acts
    as the final safety net. Hard-coding utf-8 (the old behaviour) silently
    corrupted cp1252/latin-1 files into � characters.
    """
    if not data:
        return ""
    # 1. UTF-8 (utf-8-sig also consumes a leading BOM). Strict so non-UTF-8
    #    bytes fall through to the legacy attempts instead of mojibake.
    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError:
        pass
    # 2. Windows-1252 (strict): the deterministic, correct answer for the common
    #    legacy-office case. Its 5 undefined bytes raise and fall through.
    try:
        return data.decode("cp1252").lstrip(_FEFF)
    except UnicodeDecodeError:
        pass
    # 3. charset-normalizer (optional) — only reached for cp1252-invalid bytes,
    #    where statistical detection of the real codepage beats a blind fallback.
    try:
        from charset_normalizer import from_bytes  # optional dependency

        best = from_bytes(data).best()
        if best is not None:
            return str(best).lstrip(_FEFF)
    except Exception:  # missing dep or detection failure — degrade gracefully
        pass
    # 4. latin-1 is total (all 256 bytes map) so this is the guaranteed decode.
    try:
        return data.decode("latin-1").lstrip(_FEFF)
    except UnicodeDecodeError:
        # 5. Truly pathological — never lose the read; replace the bad bytes.
        return data.decode("utf-8", errors="replace").lstrip(_FEFF)


def _check_size(p: Path, label: str) -> None:
    """Raise a clear error if ``p`` exceeds the whole-file read ceiling."""
    size = p.stat().st_size
    if size > _MAX_BYTES:
        mb = size / (1024 * 1024)
        raise ValueError(
            f"{label} is too large to read safely ({mb:.0f} MB > 100 MB) — "
            "extract a slice (page_range/sheet) or use a smaller export"
        )


def _read_unknown(p: Path) -> str:
    """Read an unknown-suffix file as text, streaming-limited and NUL-sniffed."""
    _check_size(p, "file")
    # Sniff the head for a NUL byte (a reliable binary marker) BEFORE committing
    # to decode the whole file, and read in bounded chunks rather than one big
    # read_bytes() of a potentially huge unknown blob.
    chunks: list[bytes] = []
    with open(p, "rb") as f:
        head = f.read(65536)
        if b"\x00" in head:
            raise ValueError(f"unsupported binary file type: {p.suffix or p.name!r}")
        chunks.append(head)
        while True:
            block = f.read(1024 * 1024)
            if not block:
                break
            if b"\x00" in block:
                raise ValueError(
                    f"unsupported binary file type: {p.suffix or p.name!r}"
                )
            chunks.append(block)
    return _decode_bytes(b"".join(chunks))


def _guard_office_encrypted(p: Path) -> None:
    """Raise a clear error if an OOXML file is really an encrypted OLE container.

    A password-protected .docx/.xlsx/.pptx is not a zip — it is an OLE/CFB
    "EncryptedPackage". python-docx/openpyxl would otherwise surface the
    misleading ``BadZipFile: File is not a zip file``.
    """
    try:
        with open(p, "rb") as f:
            head = f.read(8)
    except OSError:
        return
    if head.startswith(_OLE_MAGIC):
        kind = p.suffix.lstrip(".") or "office"
        raise ValueError(
            f"{kind} file is password-protected/encrypted — provide an "
            "unprotected copy to read it"
        )


def _nonempty_or_sentinel(text: str, *, kind: str) -> str:
    """Return ``text``, or a clear sentinel when it is effectively empty.

    Lets an agent tell "unreadable (scanned/image-only)" apart from a document
    that legitimately contains no text, instead of an ambiguous empty string.
    """
    if text.strip():
        return text
    if kind == "pdf":
        return _SCANNED_PDF_SENTINEL
    return (
        f"[no extractable text in this {kind} — it may be image-only, "
        "empty, or password-protected]"
    )


def _page_indices(spec: str | None, total: int) -> Iterable[int]:
    """Parse a 1-based inclusive page/slide spec into sorted 0-based indices.

    Accepts single pages, ``a-b`` ranges (open-ended ``a-`` / ``-b`` allowed),
    and comma-separated lists, e.g. ``"1,3-5,8-"``. Out-of-range numbers are
    clamped away. ``None`` means "everything".
    """
    if spec is None or str(spec).strip() == "":
        return range(total)
    keep: set[int] = set()
    for raw in str(spec).split(","):
        part = raw.strip()
        if not part:
            continue
        try:
            if "-" in part:
                a, _, b = part.partition("-")
                start = int(a) if a.strip() else 1
                end = int(b) if b.strip() else total
            else:
                start = end = int(part)
        except ValueError:
            raise ValueError(f"invalid page_range {spec!r}") from None
        for n in range(start, end + 1):
            if 1 <= n <= total:
                keep.add(n - 1)
    return sorted(keep)


# --- structured-format readers ------------------------------------------------


def _read_pdf(p: Path, *, page_range: str | None = None) -> str:
    _check_size(p, "PDF")
    from pypdf import PdfReader

    reader = PdfReader(str(p))
    if reader.is_encrypted:
        # Try the empty password (common for owner-only restrictions). A return
        # of 0 / PasswordType.NOT_DECRYPTED means it is still locked.
        try:
            unlocked = reader.decrypt("")
        except Exception:
            unlocked = 0
        if not unlocked:
            raise ValueError(
                "PDF is password-protected — provide the password or an "
                "unlocked copy to read it"
            )
    pages = reader.pages
    idxs = _page_indices(page_range, len(pages))
    text = "\n".join((pages[i].extract_text() or "") for i in idxs)
    return _nonempty_or_sentinel(text, kind="pdf")


def _iter_block_items(doc):
    """Yield the document body's paragraphs and tables IN DOCUMENT ORDER.

    ``doc.paragraphs`` and ``doc.tables`` are separate flat lists that lose the
    interleaving, so a table between two paragraphs would float out of place.
    Walking the body XML keeps tables inline with the surrounding prose.
    """
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, doc)
        elif child.tag == qn("w:tbl"):
            yield Table(child, doc)


def _docx_notes(doc) -> list[str]:
    """Extract footnote/endnote text (python-docx has no first-class API).

    The notes live in the ``word/footnotes.xml`` / ``word/endnotes.xml`` parts;
    parse their ``w:t`` runs directly. Default separator entries carry no text
    and are naturally skipped by the strip check.
    """
    W_T = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"
    notes: list[str] = []
    try:
        from lxml import etree

        parts = list(doc.part.package.iter_parts())
    except Exception:
        return notes
    for part in parts:
        name = str(getattr(part, "partname", ""))
        if not (name.endswith("footnotes.xml") or name.endswith("endnotes.xml")):
            continue
        try:
            root = etree.fromstring(part.blob)
        except Exception:
            continue
        for t in root.iter(W_T):
            if t.text and t.text.strip():
                notes.append(t.text)
    return notes


def _read_docx(p: Path) -> str:
    _guard_office_encrypted(p)
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(str(p))
    parts: list[str] = []

    # Section headers first — letterheads, document titles and running banners
    # (often the single highest-value text) live here and were dropped before.
    for section in doc.sections:
        header = section.header
        if header is not None and not header.is_linked_to_previous:
            for para in header.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

    # Body in document order so tables stay inline with surrounding paragraphs.
    for block in _iter_block_items(doc):
        if isinstance(block, Paragraph):
            parts.append(block.text)
        elif isinstance(block, Table):
            for row in block.rows:
                parts.append("\t".join(cell.text for cell in row.cells))

    # Section footers (page footers, disclaimers, addresses).
    for section in doc.sections:
        footer = section.footer
        if footer is not None and not footer.is_linked_to_previous:
            for para in footer.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

    parts.extend(_docx_notes(doc))
    return _nonempty_or_sentinel("\n".join(parts), kind="docx")


def _fmt_has_time(number_format: str | None) -> bool:
    """True if an Excel number-format code renders a time component.

    ``h``/``s``/``AM/PM``/``A/P`` are unambiguous time tokens; ``m`` alone is
    skipped because it is month-or-minute depending on neighbours.
    """
    if not number_format:
        return False
    fmt = number_format.lower()
    return any(tok in fmt for tok in ("h", "s", "am/pm", "a/p"))


def _fmt_cell(value, number_format: str | None) -> str:
    """Render a cell value as text, formatting dates/times by TYPE.

    A date-only value (midnight and no time in its number-format) becomes a bare
    ISO date — no spurious ``" 00:00:00"`` — while a value that actually carries
    a time keeps it.
    """
    if value is None:
        return ""
    # datetime is a subclass of date, so test it FIRST.
    if isinstance(value, _dt.datetime):
        midnight = value.hour == 0 and value.minute == 0 and value.second == 0
        if midnight and not _fmt_has_time(number_format):
            return value.date().isoformat()
        return value.isoformat(sep=" ")
    if isinstance(value, _dt.date):
        return value.isoformat()
    if isinstance(value, _dt.time):
        return value.isoformat()
    return str(value)


def _read_xlsx(p: Path, *, sheet: str | int | None = None) -> str:
    _guard_office_encrypted(p)
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(p), read_only=True, data_only=True)
    # Second (lazy) pass with cached values OFF: when data_only yields None for a
    # formula cell (the workbook was never opened in Excel, so no cached value),
    # fall back to the formula TEXT so computed columns aren't silently blank.
    formula_wb = None
    try:
        targets = _select_sheets(wb, sheet)
        parts: list[str] = []
        for ws in targets:
            parts.append(f"## {ws.title}")
            fws = None
            for row in ws.iter_rows():
                cells: list[str] = []
                need_formula = False
                rendered: list[str] = []
                for cell in row:
                    val = _fmt_cell(
                        cell.value, getattr(cell, "number_format", None)
                    )
                    rendered.append(val)
                    if val == "" and cell.value is None:
                        need_formula = True
                    cells.append(cell)
                if need_formula:
                    if formula_wb is None:
                        formula_wb = load_workbook(
                            filename=str(p), read_only=True, data_only=False
                        )
                    if fws is None:
                        fws = formula_wb[ws.title]
                    # Zip the same row from the formula pass; show any "=..." text
                    # where the value pass had nothing.
                    frow = _row_at(fws, cells[0].row) if cells else None
                    if frow is not None:
                        for i, fcell in enumerate(frow):
                            if i < len(rendered) and rendered[i] == "":
                                fval = fcell.value
                                if isinstance(fval, str) and fval.startswith("="):
                                    rendered[i] = fval
                parts.append("\t".join(rendered))
        return "\n".join(parts)
    finally:
        wb.close()
        if formula_wb is not None:
            formula_wb.close()


def _row_at(ws, row_number: int):
    """Return the cells of ``ws`` at 1-based ``row_number`` (read-only safe)."""
    for row in ws.iter_rows(min_row=row_number, max_row=row_number):
        return row
    return None


def _select_sheets(wb, sheet: str | int | None):
    """Resolve ``sheet`` (name or 0-based index) to a list of worksheets."""
    if sheet is None:
        return list(wb.worksheets)
    names = wb.sheetnames
    if isinstance(sheet, int) or (isinstance(sheet, str) and sheet.strip().lstrip("-").isdigit()):
        idx = int(sheet)
        if 0 <= idx < len(names):
            return [wb[names[idx]]]
    if isinstance(sheet, str) and sheet in names:
        return [wb[sheet]]
    raise ValueError(
        f"sheet {sheet!r} not found — available sheets: {', '.join(names)}"
    )


def _collect_pptx_shapes(shapes, out: list[str]) -> None:
    """Append text from shapes, recursing groups and emitting table cells."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_pptx_shapes(shape.shapes, out)  # recurse group members
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                out.append("\t".join(cell.text for cell in row.cells))
            continue
        if shape.has_text_frame and shape.text:
            out.append(shape.text)


def _read_pptx(p: Path, *, page_range: str | None = None) -> str:
    _guard_office_encrypted(p)
    import pptx

    prs = pptx.Presentation(str(p))
    slides = list(prs.slides)
    idxs = _page_indices(page_range, len(slides))
    parts: list[str] = []
    for i in idxs:
        slide = slides[i]
        texts: list[str] = []
        _collect_pptx_shapes(slide.shapes, texts)
        # Speaker notes carry the real narrative of many decks — include them.
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text
            if note and note.strip():
                texts.append(note)
        parts.append("\n".join(texts))
    return _nonempty_or_sentinel("\n".join(parts), kind="pptx")


def _read_csv(p: Path) -> str:
    # Encoding-detected: Excel/Windows CSVs are frequently UTF-8-BOM or cp1252,
    # which the old hard-coded utf-8 read corrupted into replacement chars.
    text = _decode_bytes(p.read_bytes())
    rows: list[str] = ["\t".join(row) for row in csv.reader(io.StringIO(text))]
    return "\n".join(rows)


def _read_rtf(p: Path) -> str:
    raw = _decode_bytes(p.read_bytes())
    try:
        from striprtf.striprtf import rtf_to_text  # optional, best fidelity

        return rtf_to_text(raw)
    except Exception:
        # Fallback: strip control words / groups so the user sees the prose
        # instead of raw ``\rtf1\ansi...`` noise.
        return _strip_rtf_controls(raw)


def _strip_rtf_controls(rtf: str) -> str:
    """Minimal RTF-to-text: drop hex escapes, control words and group braces."""
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", rtf)  # \'hh escaped bytes
    text = re.sub(r"\\par[d]?", "\n", text)  # paragraph breaks -> newlines
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)  # remaining control words
    text = text.replace("{", "").replace("}", "")  # group delimiters
    # Collapse the runs of blank lines the stripping leaves behind.
    return re.sub(r"\n[ \t]*\n+", "\n\n", text).strip()


def _describe_image(p: Path) -> str:
    try:
        from PIL import Image
    except ImportError:  # Pillow not installed (e.g. trimmed from a frozen build)
        return f"[image {p.suffix.lstrip('.').upper() or 'IMAGE'}: Pillow not available for metadata]"

    with Image.open(p) as img:
        fmt = img.format or (p.suffix.lstrip(".").upper() or "IMAGE")
        width, height = img.size
        mode = img.mode
    return f"[image {fmt} {width}x{height}, mode {mode}]"
