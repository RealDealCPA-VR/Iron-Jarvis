r"""PII detection + format-preserving redaction (the ``redact_pii`` tool).

Detection is DETERMINISTIC — regex patterns for structured identifiers (SSN,
ITIN, EIN, email, phone, credit card with a Luhn check, context-gated bank
numbers and dates of birth, street addresses, IPs) plus caller-supplied
``extra_terms`` for the unstructured PII only a reader can spot (person names,
employers). No LLM in the loop here: what gets redacted is exactly what the
rules + terms say, auditable from the tool call itself.

Detection runs PER LINE (v1.154.0). The separators in these patterns include
``\s``, which matches a newline, so a value ending one line and a number
starting the next were welded into a match: on a real tax return six of seven
"phone" hits were that, and ownership percentages were being blacked out. A PII
value never spans a line break, so scoping detection costs nothing.

Redaction PRESERVES the document: docx/xlsx/pptx are rewritten in place
(styles, tables, headers/footers intact — only matched characters change),
plain-text formats are string-rewritten, and PDFs are edited IN PLACE by
``pdf_redact`` (pikepdf rewrites the content stream so the glyphs are really
deleted, pdfplumber supplies the geometry for true black boxes, and the written
file is re-read to PROVE the values are gone). A PDF whose fonts cannot be
matched, or whose output cannot be verified, falls back to the old rebuild from
extracted text — layout approximate, PII genuinely gone — and the note says
which path produced the file. The source is NEVER touched; output always lands
in a new file.

A source with NO TEXT LAYER — a scanned/photographed PDF, or an image — is
REFUSED by BOTH entry points (v1.174.0, :data:`NO_TEXT_LAYER_ERROR`): the one
that rewrites a document (:func:`redact_file`) and the one that CERTIFIES it
(:func:`scan_document`). Detection here runs over extracted text, so such a
file used to produce zero findings, an identical "redacted" copy, and the
sentence "no PII found": a certification of cleanliness for a document the tool
never read a word of. Guarding only the rewrite is not enough — nobody redacts
a file they have just been told is clean. OCR can recover the text for REVIEW
(``redact_scan`` passes it in via ``text=``), but recovered text cannot be
rewritten out of an image, so redaction refuses rather than shipping a copy
that looks safe.

Styles: ``black`` = same-length █ blocks (layout preserved), ``label`` =
``[SSN]``-style category tags, ``remove`` = deleted outright.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Callable

# --------------------------------------------------------------- detection ---

#: Category -> compiled pattern. Group 1, when present, is the PII portion
#: (context words stay); otherwise the whole match is the PII.
_PATTERNS: dict[str, re.Pattern[str]] = {
    # 123-45-6789 (also space-separated). Area 9xx is an ITIN, matched below.
    "ssn": re.compile(r"\b(?!000|666|9\d\d)\d{3}[- ](?!00)\d{2}[- ](?!0000)\d{4}\b"),
    "itin": re.compile(r"\b9\d{2}[- ]\d{2}[- ]\d{4}\b"),
    # Labeled contiguous 9-digit SSN/ITIN ("SSN: 123456789").
    "ssn_labeled": re.compile(
        r"\b(?:ssn|itin|social security(?:\s+(?:number|no\.?))?)\s*[:#]?\s*"
        r"(\d{3}[- ]?\d{2}[- ]?\d{4})",
        re.IGNORECASE,
    ),
    "ein": re.compile(r"\b\d{2}-\d{7}\b"),
    "email": re.compile(r"\b[\w.+-]+@[\w-]+(?:\.[\w-]+)+\b"),
    "phone": re.compile(r"(?:\+?1[-.\s]?)?\(?\d{3}\)?[-.\s]\d{3}[-.\s]?\d{4}\b"),
    # 13-19 digits (spaces/dashes allowed) — confirmed by Luhn below.
    "credit_card": re.compile(r"\b(?:\d[ -]?){12,18}\d\b"),
    "bank_account": re.compile(
        r"\b(?:account|acct|routing|aba)\s*(?:number|no\.?|#)?\s*[:#]?\s*(\d{6,17})\b",
        re.IGNORECASE,
    ),
    "dob": re.compile(
        r"\b(?:dob|date of birth|born(?:\s+on)?|birth\s*date)\s*[:#]?\s*"
        r"((?:\d{1,2}[/-]\d{1,2}[/-]\d{2,4})|(?:[A-Z][a-z]+\.?\s+\d{1,2},?\s+\d{4}))",
        re.IGNORECASE,
    ),
    "address": re.compile(
        r"\b\d{1,6}\s+(?:[A-Z][\w'.-]*\s){0,4}?"
        r"(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Drive|Dr|"
        r"Court|Ct|Circle|Cir|Way|Place|Pl|Terrace|Ter|Highway|Hwy|Parkway|"
        r"Pkwy|Trail|Trl|Loop)\.?\b"
        r"(?:\s*,?\s*(?:#|Apt\.?|Suite|Ste\.?|Unit)\s*\w+)?",
        # CASE-INSENSITIVE (v1.154.0). Without this the street suffixes only
        # matched mixed case, and TAX DOCUMENTS ARE UPPERCASE: "5059 ALAMANDA
        # DR" on a real K-1 was not detected, so a client's home address stayed
        # in a file the user believed was redacted. The trade is a few more
        # candidates ("12 St" in prose), and it is the right way round — this
        # tool is confirm-first, so a false positive costs one glance while a
        # miss leaks a home address.
        re.IGNORECASE,
    ),
    "ip": re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b"),
}

#: Short display tags for the ``label`` style.
_LABELS: dict[str, str] = {
    "ssn": "SSN",
    "itin": "ITIN",
    "ssn_labeled": "SSN",
    "ein": "EIN",
    "email": "EMAIL",
    "phone": "PHONE",
    "credit_card": "CARD",
    "bank_account": "ACCOUNT",
    "dob": "DOB",
    "address": "ADDRESS",
    "ip": "IP",
    "custom": "REDACTED",
}

#: Display labels for the REVIEW list, which has different needs from the
#: in-document replacement text. "custom" renders as "[REDACTED]" inside a
#: document — correct there — but as a badge next to a value the user typed
#: themselves it reads as a status ("this was redacted"), not as what it is:
#: the one candidate that came from THEM rather than from a pattern.
_SCAN_LABELS: dict[str, str] = {**_LABELS, "custom": "TERM"}

ALL_CATEGORIES: frozenset[str] = frozenset(_LABELS)

#: Redaction styles the tool accepts.
STYLES = ("black", "label", "remove")


def _luhn_ok(digits: str) -> bool:
    ds = [int(c) for c in digits if c.isdigit()]
    if not 13 <= len(ds) <= 19:
        return False
    total = 0
    for i, d in enumerate(reversed(ds)):
        if i % 2 == 1:
            d *= 2
            if d > 9:
                d -= 9
        total += d
    return total % 10 == 0


def _categorize_term(term: str) -> str:
    """The category a CONFIRMED literal term belongs to (first pattern that
    matches the whole term wins), so 'label'-style redaction still tags it
    ``[SSN]`` rather than a generic ``[REDACTED]``."""
    for cat, rx in _PATTERNS.items():
        m = rx.fullmatch(term) or rx.match(term)
        if m is not None and m.group(0) == term:
            if cat == "credit_card" and not _luhn_ok(term):
                continue
            return cat
    return "custom"


def _iter_lines(text: str):
    """Yield ``(offset, line)`` for each line, offsets into the ORIGINAL text.

    Used to scope pattern detection to a single line — see the note in
    :func:`find_pii_spans`. Keeps the caller's span coordinates absolute, so
    nothing downstream has to know detection was chunked.
    """
    offset = 0
    for line in text.splitlines(keepends=True):
        yield offset, line
        offset += len(line)


def find_pii_spans(
    text: str,
    *,
    extra_terms: list[str] | None = None,
    categories: set[str] | frozenset[str] | None = None,
    only_terms: list[str] | None = None,
) -> list[tuple[int, int, str]]:
    """Return non-overlapping ``(start, end, category)`` spans, sorted by start.
    ``extra_terms`` are matched literally (case-insensitive) as ``custom``.
    ``only_terms`` switches to CONFIRMED mode: pattern detection is skipped
    entirely and exactly those literal values are redacted (case-insensitive),
    each categorized for labeling — the human-in-the-loop contract where what
    gets removed is precisely what was confirmed, nothing else.
    Earlier-starting/longer spans win overlaps."""
    raw: list[tuple[int, int, str]] = []
    if only_terms is not None:
        for term in only_terms:
            t = (term or "").strip()
            if len(t) < 2:
                continue  # a 1-char term would shred the document
            cat = _categorize_term(t)
            for m in re.finditer(re.escape(t), text, re.IGNORECASE):
                raw.append((m.start(), m.end(), cat))
    else:
        wanted = set(categories) if categories else set(_PATTERNS) | {"custom"}
        # PER LINE, never across a line break (v1.154.0). The separators in
        # these patterns include ``\s``, which matches a NEWLINE, so a value
        # ending one line and a number starting the next were being welded into
        # a match. On the tax return that prompted this, SIX of seven "phone"
        # hits were that: '1096\n100.0000' and '416\n100.0000' are ownership
        # percentages, and the redactor was blacking out real financial figures
        # while the user had no way to see why. A PII value never spans a line
        # break in a real document, so scoping detection to one line costs
        # nothing and removes the whole class of false positive.
        for line_start, line in _iter_lines(text):
            for cat, rx in _PATTERNS.items():
                if cat not in wanted:
                    continue
                for m in rx.finditer(line):
                    start, end = (m.span(1) if m.groups() and m.group(1) else m.span())
                    value = line[start:end]
                    if cat == "credit_card" and not _luhn_ok(value):
                        continue
                    if cat == "ip" and any(
                        int(p) > 255 for p in re.findall(r"\d+", value)
                    ):
                        continue
                    raw.append((line_start + start, line_start + end, cat))
        if "custom" in wanted:
            for term in extra_terms or []:
                t = (term or "").strip()
                if len(t) < 2:
                    continue  # a 1-char term would shred the document
                for m in re.finditer(re.escape(t), text, re.IGNORECASE):
                    raw.append((m.start(), m.end(), "custom"))
    # Resolve overlaps: sort by (start, -length); keep spans that don't overlap
    # an already-kept one.
    raw.sort(key=lambda s: (s[0], -(s[1] - s[0])))
    kept: list[tuple[int, int, str]] = []
    last_end = -1
    for start, end, cat in raw:
        if start >= last_end:
            kept.append((start, end, cat))
            last_end = end
    return kept


def _replacement(value: str, category: str, style: str) -> str:
    if style == "remove":
        return ""
    if style == "label":
        return f"[{_LABELS.get(category, 'REDACTED')}]"
    return "█" * len(value)  # "black" — same length keeps layout intact


#: A resolved span ready to apply: (start, end, category, replacement).
_Span = tuple[int, int, str, str]


def _make_spans_fn(
    style: str,
    extra_terms: list[str] | None,
    categories: set[str] | frozenset[str] | None,
    only_terms: list[str] | None = None,
) -> Callable[[str], list[_Span]]:
    def spans_for(text: str) -> list[_Span]:
        found = find_pii_spans(
            text, extra_terms=extra_terms, categories=categories,
            only_terms=only_terms,
        )
        return [(s, e, cat, _replacement(text[s:e], cat, style)) for s, e, cat in found]

    return spans_for


def _apply_spans(text: str, spans: list[_Span]) -> tuple[str, dict[str, int]]:
    counts: dict[str, int] = {}
    out: list[str] = []
    cursor = 0
    for start, end, cat, repl in spans:
        out.append(text[cursor:start])
        out.append(repl)
        counts[cat] = counts.get(cat, 0) + 1
        cursor = end
    out.append(text[cursor:])
    return "".join(out), counts


def mask_text(
    text: str,
    *,
    style: str = "black",
    extra_terms: list[str] | None = None,
    categories: set[str] | frozenset[str] | None = None,
    only_terms: list[str] | None = None,
) -> tuple[str, dict[str, int]]:
    """Redact *text*; returns ``(redacted, counts_by_category)``."""
    spans = _make_spans_fn(style, extra_terms, categories, only_terms)(text)
    return _apply_spans(text, spans)


# ----------------------------------------------------------------- scanning ---

#: Context chars shown either side of a finding in scan results.
_SCAN_CONTEXT = 40
#: Findings cap — a pathological document must not flood the review.
_MAX_FINDINGS = 120


def scan_text(
    text: str,
    *,
    extra_terms: list[str] | None = None,
    categories: set[str] | frozenset[str] | None = None,
) -> list[dict[str, Any]]:
    """STEP 1 of confirmed redaction: every distinct PII candidate in *text*,
    grouped by (category, value) with an occurrence count and one context
    snippet — the reviewable list a human confirms before anything is
    removed. Ordered by first appearance; capped with an honest flag."""
    spans = find_pii_spans(text, extra_terms=extra_terms, categories=categories)
    findings: dict[tuple[str, str], dict[str, Any]] = {}
    for start, end, cat in spans:
        value = text[start:end]
        key = (cat, value.lower())
        if key in findings:
            findings[key]["count"] += 1
            continue
        ctx_start = max(0, start - _SCAN_CONTEXT)
        ctx_end = min(len(text), end + _SCAN_CONTEXT)
        context = " ".join(text[ctx_start:ctx_end].split())
        findings[key] = {
            "category": cat,
            "label": _SCAN_LABELS.get(cat, "REDACTED"),
            "value": value,
            "count": 1,
            "context": context,
        }
        if len(findings) >= _MAX_FINDINGS:
            break
    out = list(findings.values())
    for i, f in enumerate(out, start=1):
        f["id"] = i
    return out


def scan_document(
    path: Path,
    *,
    extra_terms: list[str] | None = None,
    categories: set[str] | frozenset[str] | None = None,
    text: str | None = None,
) -> list[dict[str, Any]]:
    """Scan a document file for PII candidates (see :func:`scan_text`).

    ``text`` (v1.174.0) supplies text the caller ALREADY has — specifically an
    OCR transcription of a scanned document, which has no text layer for
    ``extract_text`` to find. Omitted, this reads the file itself — and then
    REFUSES a file whose words are pixels (:data:`NO_TEXT_LAYER_ERROR`).

    That refusal is the whole point of the function's safety story, and it was
    missing from the leg that matters most. ``POST /documents/redact/scan``
    calls this with no ``text``: on a scanned tax return the patterns ran over
    the "no extractable text" sentinel, found nothing, and the Documents page
    rendered a green shield reading "No personal data found." for a document
    the app had never read a word of. The APPLY leg refusing is not enough —
    a user told the file is clean never reaches apply. A file some earlier OCR
    already transcribed is NOT refused: ``extract_text`` serves that
    transcription (it carries the OCR marker), so the candidates below are the
    ones really in the document.
    """
    from .readers import extract_text

    if text is None:
        text = extract_text(path)
        _refuse_if_no_text_layer(Path(path), text)
    return scan_text(text, extra_terms=extra_terms, categories=categories)


# -------------------------------------------------------- format redactors ---


def _merge_counts(total: dict[str, int], part: dict[str, int]) -> None:
    for k, v in part.items():
        total[k] = total.get(k, 0) + v


def _redact_runs(
    runs: list[Any], spans_for: Callable[[str], list[_Span]]
) -> dict[str, int]:
    """Redact PII across a paragraph's runs, PRESERVING run formatting.

    Matches are found on the CONCATENATED text (PII often spans runs — e.g.
    a bold SSN split by the editor), then each run's slice is rewritten. A
    replacement whose length differs (label/remove styles) lands wholly in the
    run where the match STARTS; later runs' matched characters are dropped.
    """
    texts = [r.text or "" for r in runs]
    combined = "".join(texts)
    if not combined:
        return {}
    spans = spans_for(combined)
    if not spans:
        return {}
    counts: dict[str, int] = {}
    for _s, _e, cat, _r in spans:
        counts[cat] = counts.get(cat, 0) + 1
    offsets: list[int] = []
    pos = 0
    for t in texts:
        offsets.append(pos)
        pos += len(t)
    span_iter = iter(spans)
    span = next(span_iter, None)
    for i, t in enumerate(texts):
        rs, re_ = offsets[i], offsets[i] + len(t)
        cursor = rs
        parts: list[str] = []
        while cursor < re_:
            if span is None or span[0] >= re_:
                parts.append(combined[cursor:re_])
                break
            s, e, _cat, repl = span
            if s > cursor:
                parts.append(combined[cursor:s])
                cursor = s
                continue
            # Inside the span: the replacement is emitted only by the run where
            # the span STARTS; runs it merely continues into contribute nothing.
            if s >= rs:
                parts.append(repl)
            cursor = min(e, re_)
            if e <= re_:
                span = next(span_iter, None)
        new_text = "".join(parts)
        if (t or "") != new_text:
            runs[i].text = new_text
    return counts


def _redact_docx(src: Path, dst: Path, spans_for) -> dict[str, int]:
    import docx  # python-docx

    doc = docx.Document(str(src))
    counts: dict[str, int] = {}

    def do_paragraphs(paragraphs) -> None:
        for par in paragraphs:
            _merge_counts(counts, _redact_runs(list(par.runs), spans_for))

    def do_tables(tables) -> None:
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    do_paragraphs(cell.paragraphs)
                    do_tables(cell.tables)  # nested tables

    do_paragraphs(doc.paragraphs)
    do_tables(doc.tables)
    for section in doc.sections:
        for part in (section.header, section.footer):
            do_paragraphs(part.paragraphs)
            do_tables(part.tables)
    doc.save(str(dst))
    return counts


def _redact_xlsx(src: Path, dst: Path, spans_for) -> dict[str, int]:
    from openpyxl import load_workbook

    wb = load_workbook(str(src))  # formulas preserved (not data_only)
    counts: dict[str, int] = {}
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                # Only string cells: numbers/formulas/dates stay untouched (a
                # formula rewrite could corrupt the sheet; noted in the tool).
                if isinstance(v, str) and not v.startswith("="):
                    masked, part = _apply_spans(v, spans_for(v))
                    if part:
                        cell.value = masked
                        _merge_counts(counts, part)
    wb.save(str(dst))
    return counts


def _redact_pptx(src: Path, dst: Path, spans_for) -> dict[str, int]:
    from pptx import Presentation

    prs = Presentation(str(src))
    counts: dict[str, int] = {}

    def do_text_frame(tf) -> None:
        for par in tf.paragraphs:
            _merge_counts(counts, _redact_runs(list(par.runs), spans_for))

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                do_text_frame(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        do_text_frame(cell.text_frame)
    prs.save(str(dst))
    return counts


#: Raised (as a ValueError) for a source whose words are PIXELS. Every redactor
#: below rewrites TEXT, so on a scanned PDF the old path found no spans, wrote
#: an identical copy, and reported "no PII found" — a tax return full of SSNs
#: certified clean by a tool that never read a word of it. Refusing is the only
#: honest answer: there is no rewrite of a text layer that does not exist.
NO_TEXT_LAYER_ERROR = (
    "this document has NO TEXT LAYER (it is a scan/photo/image) — nothing was "
    "scanned for PII and nothing can be redacted by rewriting text. A rebuilt "
    "copy would look redacted and still carry every value, and an empty finding "
    "list here means NOT SCANNED, never CLEAN. Transcribe or re-scan it with a "
    "text layer, or redact the image itself"
)


def _refuse_if_no_text_layer(src: Path, text: str) -> None:
    """Guard EVERY entry point that reads a file here — :func:`scan_document`
    (the leg that CERTIFIES a document clean) and :func:`redact_file` (the leg
    that rewrites one), tool and route alike.

    ``needs_ocr`` rather than ``looks_scanned_pdf``: a ``.png``/``.jpg`` of a
    W-2 has no text layer either and would otherwise be scanned for PII as the
    string ``"[image PNG 900x1200, mode RGB]"`` — clean by construction. It
    also recognises a transcript we already produced (the OCR marker), so a
    file whose text was recovered earlier is scanned, not refused.
    """
    from .ocr import needs_ocr

    if needs_ocr(src, text):
        raise ValueError(NO_TEXT_LAYER_ERROR)


def _redact_pdf(src: Path, dst: Path, spans_for) -> tuple[dict[str, int], str]:
    """Redact a PDF IN PLACE when we can prove it worked; rebuild when we can't.

    The in-place path (``pdf_redact``, v1.154.0) keeps the real page — form
    rules, page size, fonts — and deletes the PII glyphs from the content
    stream, then RE-READS the written file to prove the values are gone. The
    rebuild below is what shipped before and stays as the fallback, because
    "truly gone, layout approximate" still beats "looks right, still leaks".

    The order matters and is the whole point: try to keep the document, but
    never keep it at the cost of the guarantee.
    """
    from .readers import extract_text
    from .writers import write_document

    # The file's REAL text layer: `use_ocr_cache=False` because a cached
    # transcription proves the words are there, it does not make them
    # rewritable — rebuilding a scan from OCR text would destroy the document
    # AND leave the original pixels in whatever the user shares next.
    text = extract_text(src, use_ocr_cache=False)
    _refuse_if_no_text_layer(src, text)
    spans = spans_for(text)
    masked, counts = _apply_spans(text, spans)

    # The exact strings the caller's rules matched — the same values the
    # in-place path must remove and then fail to find.
    values = sorted({text[s:e] for s, e, _cat, _repl in spans if e > s}, key=len, reverse=True)
    if values:
        try:
            from .pdf_redact import RedactionUnverified, UnsupportedPdf, redact_pdf

            replacements = {text[s:e]: repl for s, e, _cat, repl in spans if e > s}
            redact_pdf(
                src,
                dst,
                values=values,
                replacement=lambda v: replacements.get(v, " " * len(v)),
            )
            return counts, (
                "redacted in place — the original pages, form lines and fonts "
                "are untouched, and the removed text is verified gone from the "
                "output (not merely covered)."
            )
        except (UnsupportedPdf, RedactionUnverified, ImportError) as exc:
            reason = str(exc)
        except Exception as exc:  # noqa: BLE001 — any surprise falls back too
            reason = f"{type(exc).__name__}: {exc}"
    else:
        reason = ""

    write_document(dst, masked, kind="pdf")
    note = (
        "PDF REBUILT from extracted text — the content is truly removed but the "
        "layout is approximate (page size, form lines and fonts are not "
        "preserved)."
        + (f" In-place redaction was not usable here: {reason}" if reason else "")
    )
    return counts, note


_TEXT_SUFFIXES = {
    ".txt", ".md", ".csv", ".tsv", ".html", ".htm", ".json", ".log", ".xml",
    ".yaml", ".yml", ".rtf",
}


def redact_file(
    src: Path,
    dst: Path,
    *,
    style: str = "black",
    extra_terms: list[str] | None = None,
    categories: set[str] | frozenset[str] | None = None,
    only_terms: list[str] | None = None,
) -> tuple[dict[str, int], str]:
    """Redact *src* into *dst* (same format). Returns ``(counts, note)``.
    ``only_terms`` = CONFIRMED mode: exactly those values are redacted and
    nothing else (see :func:`find_pii_spans`). The source is never modified."""
    if style not in STYLES:
        raise ValueError(f"unknown style: {style!r} (use black, label, or remove)")
    spans_for = _make_spans_fn(style, extra_terms, categories, only_terms)
    suffix = src.suffix.lower()
    note = ""
    if suffix == ".docx":
        counts = _redact_docx(src, dst, spans_for)
    elif suffix == ".xlsx":
        counts = _redact_xlsx(src, dst, spans_for)
        note = "string cells redacted; numeric cells and formulas are untouched"
    elif suffix == ".pptx":
        counts = _redact_pptx(src, dst, spans_for)
    elif suffix == ".pdf":
        counts, note = _redact_pdf(src, dst, spans_for)
    elif suffix in _TEXT_SUFFIXES or suffix == "":
        # Decode the way every other reader in this package does (utf-8-sig →
        # strict cp1252 → charset-normalizer → latin-1) instead of hard-coding
        # utf-8 with errors="replace". That old read turned EVERY non-UTF-8
        # character of a legacy Windows/office export into U+FFFD and then wrote
        # it into the deliverable, so a copy whose only advertised change was
        # the PII differed from the original everywhere — and `read_document`
        # on the same file rendered it correctly, because it decodes properly.
        # Detection is unaffected (the patterns are ASCII), so this was purely
        # collateral corruption of a file the user then shares.
        from .readers import _decode_bytes

        raw = src.read_bytes()
        text = _decode_bytes(raw)
        masked, counts = _apply_spans(text, spans_for(text))
        # newline="" — the bytes decode carries the source's own CRLF/CR through
        # verbatim, and the default translation would re-expand each "\n" into
        # os.linesep (CR CR LF on Windows). Keep a UTF-8 BOM if the source had
        # one: Excel opens a BOM-less CSV in the legacy codepage, so dropping it
        # would mojibake the very file this fix exists to keep intact.
        enc = "utf-8-sig" if raw.startswith(b"\xef\xbb\xbf") else "utf-8"
        dst.write_text(masked, encoding=enc, newline="")
    else:
        raise ValueError(
            f"unsupported format for redaction: {suffix or '(no extension)'} — "
            "supported: .docx .xlsx .pptx .pdf and text formats "
            "(.txt .md .csv .tsv .html .json …)"
        )
    return counts, note
