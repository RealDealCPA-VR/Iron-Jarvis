"""Document writers.

``write_document(path, content, *, kind=None)`` creates a real file on disk,
choosing the format from the path suffix (or ``kind``, which overrides it).

String content is markdown-aware: it is parsed by
:mod:`iron_jarvis.documents.markdown` into blocks (headings, bullets, numbered
lists, code fences, pipe tables, ``---`` rules, and inline
``**bold**``/``*italic*``/`` `code` ``/``[link](url)``/``![img](url)`` runs)
which the rich writers render natively. Plain text with no markers simply
becomes paragraphs, so flat strings keep working everywhere.

* ``.docx`` -> python-docx: real Heading/List styles, real tables, shaded
  monospace code blocks, bold/italic/code/hyperlink runs. Non-string content:
  one paragraph per line.
* ``.xlsx`` -> openpyxl: a ``list[list]`` writes rows; a str writes one cell per
  line in column A; ``{"sheets": {name: rows}}`` writes one worksheet per key
  with a bolded/frozen header row, sized columns, and numeric/date COERCION of
  data cells. ``=...`` strings stay formulas; leading-zero ids stay text.
* ``.pptx`` -> python-pptx: each ``# `` heading or ``---`` starts a new slide;
  long sections spill onto ``Title (cont.)`` slides; tables render as real
  pptx tables; a ``Notes:`` line feeds the slide's speaker notes.
* ``.pdf``  -> fpdf2: sized bold headings, wrapped paragraphs, indented
  bullets, native wrapping tables with per-column alignment, Courier code on
  grey fill. Long tokens are pre-broken and any layout failure degrades to a
  reduced/plain render rather than producing no file.
* ``.html/.htm`` -> standalone HTML with inline CSS (links/images, validly
  nested lists, aligned tables). A full HTML page is passed through verbatim.
* ``.json`` -> ``json.dumps`` for dict/list content (str passthrough).
* ``.yaml/.yml`` -> ``yaml.safe_dump`` if PyYAML is importable, else JSON.
* ``.csv``  -> stdlib csv from a ``list[list]`` or from text lines.
* ``.txt/.md`` and anything else -> UTF-8 text.

Every writer saves to a sibling temp file then ``os.replace()`` onto the target
(atomic): a mid-save failure never clobbers an existing good file or leaves a
0-byte document. Parent directories are created as needed; the Path is returned.

BEAUTY layer (v1.134.0): ``write_document`` also takes an optional ``options``
dict — DECLARATIVE polish the engine applies so models pick options instead of
writing styling code. For ``.docx``: ``theme`` (a :mod:`.themes` name),
``cover`` (title page from the first H1; optional ``subtitle``; ``date`` True/
str/False), ``header_text``, ``footer`` (``"page-numbers"`` for a live page
number, or text — a literal ``{page}`` inside the text becomes the live field,
which is how "text plus page numbers" is expressed). For ``.xlsx``: ``theme``,
``autosize`` (defaults ON when a theme is set), ``freeze_header``, ``banded``,
``number_formats`` ({column letter: Excel format string}). An ``.xlsx`` sheet
value may also be ``{"rows": [...], "charts": [...]}`` where each chart is
``{"type": "bar"|"line"|"pie", "title", "data_range", "categories_range",
"anchor"}`` — ranges are validated (A1-style, within the used cells) and an
invalid chart is SKIPPED with a recorded warning, never a crashed write. With
no options the output is exactly the legacy output — the default flow is
untouched.
"""

from __future__ import annotations

import csv
import html as _html_mod
import json
import os
import re
from contextlib import contextmanager
from datetime import datetime
from pathlib import Path
from typing import Any

from .markdown import Block, Run, parse_markdown
from .themes import THEME_NAMES, Theme, get_theme

#: Suffixes with a dedicated writer (everything else falls back to UTF-8 text).
SUPPORTED_WRITE: set[str] = {
    ".docx",
    ".xlsx",
    ".pptx",
    ".pdf",
    ".csv",
    ".txt",
    ".md",
    ".json",
    ".html",
    ".htm",
    ".log",
    ".yaml",
    ".yml",
}


def write_document(
    path: str | Path,
    content: Any,
    *,
    kind: str | None = None,
    options: "dict[str, Any] | None" = None,
    warnings: "list[str] | None" = None,
) -> Path:
    """Write ``content`` to ``path`` as a real document. Returns the Path.

    ``options`` is the declarative beauty layer (module docstring) — honored
    by ``.docx``/``.xlsx``, ignored (with a recorded warning) elsewhere.
    ``warnings`` is a caller-supplied OUT list: non-fatal issues (skipped
    charts, unknown theme names) are appended so tools can surface them —
    a beauty problem degrades the polish, never the write.
    """
    p = Path(path)
    p.parent.mkdir(parents=True, exist_ok=True)

    suffix = ("." + kind.lstrip(".")).lower() if kind else p.suffix.lower()

    if options and suffix not in (".docx", ".xlsx"):
        # Honest, not fatal: the write proceeds exactly as without options.
        _warn(warnings, f"beauty options are honored for .docx/.xlsx only — ignored for {suffix or 'this format'}")

    if suffix == ".docx":
        _write_docx(p, content, options=options, warnings=warnings)
    elif suffix == ".xlsx":
        _write_xlsx(p, content, options=options, warnings=warnings)
    elif suffix == ".pptx":
        _write_pptx(p, content)
    elif suffix == ".pdf":
        _write_pdf(p, content)
    elif suffix == ".csv":
        _write_csv(p, content)
    elif suffix in (".html", ".htm"):
        _write_html(p, content)
    elif suffix == ".json":
        _write_json(p, content)
    elif suffix in (".yaml", ".yml"):
        _write_yaml(p, content)
    else:
        _write_text(p, content)
    return p


# --- atomic write --------------------------------------------------------------


@contextmanager
def _atomic(p: Path):
    """Yield a sibling temp path; ``os.replace`` it onto ``p`` only on success.

    If the body raises (a half-written temp), we delete the temp and re-raise,
    so the pre-existing good file at ``p`` is never touched and no 0-byte doc is
    ever left behind. The temp lives in the SAME directory so ``os.replace`` is
    an atomic same-filesystem rename on every platform (incl. Windows).
    """
    tmp = p.with_name(f".{p.name}.tmp-{os.getpid()}")
    try:
        yield tmp
        os.replace(tmp, p)
    except BaseException:
        try:
            Path(tmp).unlink()
        except OSError:
            pass
        raise


def _atomic_text(p: Path, data: str) -> None:
    with _atomic(p) as tmp:
        Path(tmp).write_text(data, encoding="utf-8")


# --- helpers ------------------------------------------------------------------


def _warn(warnings: "list[str] | None", msg: str) -> None:
    """Record a non-fatal beauty issue when the caller gave us somewhere to."""
    if warnings is not None:
        warnings.append(msg)


def _beauty_opts(
    options: "dict[str, Any] | None", warnings: "list[str] | None"
) -> "tuple[dict[str, Any] | None, Theme | None]":
    """(opts, theme) — opts is None when there is nothing to apply, which is
    the guard every beauty step hangs off (no options => the legacy path)."""
    opts = options if isinstance(options, dict) and options else None
    theme = get_theme(str(opts.get("theme") or "")) if opts else None
    if opts and opts.get("theme") and theme is None:
        _warn(
            warnings,
            f"unknown theme {opts['theme']!r} — available: {', '.join(THEME_NAMES)}",
        )
    return opts, theme


def _as_text(content: Any) -> str:
    if isinstance(content, str):
        return content
    if isinstance(content, (list, tuple)):
        return "\n".join(
            "\t".join(str(c) for c in row)
            if isinstance(row, (list, tuple))
            else str(row)
            for row in content
        )
    return str(content)


def _as_lines(content: Any) -> list[str]:
    return _as_text(content).split("\n")


def _run_parts(r: Any) -> tuple[str, bool, bool, bool, str | None, bool]:
    """(text, bold, italic, code, href, image) for a Run OR a bare 3-tuple."""
    return (
        r[0],
        r[1],
        r[2],
        getattr(r, "code", False),
        getattr(r, "href", None),
        getattr(r, "image", False),
    )


# --- json / yaml ---------------------------------------------------------------


def _write_json(p: Path, content: Any) -> None:
    # A dict/list must serialise as real JSON — the old text path wrote a Python
    # repr (single quotes, True/None) that no JSON parser accepts. Already-
    # serialised strings pass through untouched.
    if isinstance(content, str):
        data = content
    else:
        data = json.dumps(content, indent=2, ensure_ascii=False, default=str)
    _atomic_text(p, data)


def _write_yaml(p: Path, content: Any) -> None:
    if isinstance(content, str):
        _atomic_text(p, content)
        return
    try:
        import yaml  # optional dependency

        data = yaml.safe_dump(content, sort_keys=False, allow_unicode=True)
    except Exception:  # PyYAML missing or un-dumpable -> valid JSON is a fine YAML
        data = json.dumps(content, indent=2, ensure_ascii=False, default=str)
    _atomic_text(p, data)


# --- docx ----------------------------------------------------------------------


def _write_docx(
    p: Path,
    content: Any,
    options: "dict[str, Any] | None" = None,
    warnings: "list[str] | None" = None,
) -> None:
    import docx

    # BEAUTY is strictly additive: with no options this function must produce
    # exactly the document the legacy path always did, so every styling /
    # cover / header-footer step below is gated on the options dict.
    opts, theme = _beauty_opts(options, warnings)

    doc = docx.Document()
    if theme is not None:
        _docx_apply_theme(doc, theme)
    blocks = parse_markdown(content) if isinstance(content, str) else None
    if opts and opts.get("cover"):
        _docx_cover(doc, blocks, opts, theme, fallback_title=p.stem)
    if opts and (str(opts.get("header_text") or "").strip() or str(opts.get("footer") or "").strip()):
        _docx_header_footer(doc, opts)
    if blocks is not None:
        _docx_render(doc, blocks)
    else:
        for line in _as_lines(content):
            doc.add_paragraph(line)
    with _atomic(p) as tmp:
        doc.save(str(tmp))


def _docx_apply_theme(doc: Any, theme: Theme) -> None:
    """Retune the NAMED styles (Normal, Heading 1-4, Title) + page margins.

    Styling via named styles — not per-run formatting — means every paragraph
    the renderer emits (and anything the user types later in Word) inherits
    the theme, and Word's style pane still shows honest style names.
    """
    from docx.shared import Inches, Pt, RGBColor

    styles = doc.styles
    try:
        normal = styles["Normal"]
        normal.font.name = theme.body_font
        normal.font.size = Pt(theme.body_size_pt)
    except KeyError:  # pragma: no cover - Normal always exists in the template
        pass
    for level in (1, 2, 3, 4):
        try:
            st = styles[f"Heading {level}"]
        except KeyError:  # style missing from the template -> skip, never crash
            continue
        st.font.name = theme.heading_font
        st.font.size = Pt(theme.heading_size(level))
        st.font.color.rgb = RGBColor(*theme.accent_rgb)
        st.font.bold = True

    top, right, bottom, left = theme.margins_in
    for section in doc.sections:
        section.top_margin = Inches(top)
        section.right_margin = Inches(right)
        section.bottom_margin = Inches(bottom)
        section.left_margin = Inches(left)


def _docx_cover(
    doc: Any,
    blocks: "list[Block] | None",
    opts: dict[str, Any],
    theme: "Theme | None",
    *,
    fallback_title: str,
) -> None:
    """Title page: first H1 (or the filename), optional subtitle, date, then a
    page break so the body starts on page 2. The H1 stays in the body too —
    a cover ADDS a page, it never rewrites the content."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    title = next(
        (b.text for b in (blocks or []) if b.kind == "heading" and b.level == 1),
        "",
    ).strip() or str(opts.get("title") or "").strip() or fallback_title or "Document"

    for _ in range(5):  # push the title toward the visual center of the page
        doc.add_paragraph()
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = tp.add_run(title)
    run.bold = True
    run.font.size = Pt(34)
    if theme is not None:
        run.font.name = theme.heading_font
        run.font.color.rgb = RGBColor(*theme.accent_rgb)

    subtitle = str(opts.get("subtitle") or "").strip()
    if subtitle:
        sp = doc.add_paragraph()
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        srun = sp.add_run(subtitle)
        srun.italic = True
        srun.font.size = Pt(14)

    # date: True (default) -> today; a string -> verbatim; falsy -> none.
    date = opts.get("date", True)
    date_text = (
        str(date).strip()
        if isinstance(date, str)
        else (datetime.now().strftime("%B %d, %Y") if date else "")
    )
    if date_text:
        dp = doc.add_paragraph()
        dp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        dp.add_run(date_text).font.size = Pt(11)

    doc.add_page_break()


def _docx_page_number_field(paragraph: Any) -> None:
    """Append a live PAGE number field to ``paragraph``.

    python-docx has no field API, so this is the standard fieldcode dance:
    a run carrying the raw ``w:fldChar begin`` / ``w:instrText PAGE`` /
    ``w:fldChar end`` triplet, which Word renders as the current page number.
    """
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.append(begin)
    run._r.append(instr)
    run._r.append(end)


def _docx_header_footer(doc: Any, opts: dict[str, Any]) -> None:
    """Section header text + footer. ``footer`` contract: the sentinel
    ``"page-numbers"`` renders just the centered live page number; any other
    text renders as-is, with a literal ``{page}`` token replaced by the live
    field — that token is how "text AND page numbers" is expressed."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    header_text = str(opts.get("header_text") or "").strip()
    footer = str(opts.get("footer") or "").strip()
    for section in doc.sections:
        if header_text:
            hp = section.header.paragraphs[0]
            hp.text = header_text
            hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if footer:
            fp = section.footer.paragraphs[0]
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if footer == "page-numbers":
                _docx_page_number_field(fp)
            else:
                parts = footer.split("{page}")
                for i, part in enumerate(parts):
                    if i:
                        _docx_page_number_field(fp)
                    if part:
                        fp.add_run(part)


def _docx_hyperlink(paragraph: Any, url: str, text: str, *, code: bool = False) -> None:
    """Append a real clickable ``w:hyperlink`` run (blue + underlined)."""
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    r_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    link = OxmlElement("w:hyperlink")
    link.set(qn("r:id"), r_id)
    run = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "0563C1")
    rpr.append(color)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    rpr.append(underline)
    if code:
        rfonts = OxmlElement("w:rFonts")
        rfonts.set(qn("w:ascii"), "Consolas")
        rfonts.set(qn("w:hAnsi"), "Consolas")
        rpr.append(rfonts)
    run.append(rpr)
    t = OxmlElement("w:t")
    t.set(qn("xml:space"), "preserve")
    t.text = text
    run.append(t)
    link.append(run)
    paragraph._p.append(link)


def _docx_runs(paragraph: Any, runs: list[Run]) -> None:
    for r in runs:
        text, bold, italic, code, href, image = _run_parts(r)
        if href:
            # Links and images both become clickable text — the alt/label never
            # leaks the raw ``[..](..)`` markup into the document.
            label = text if not image else (text or href)
            _docx_hyperlink(paragraph, href, label, code=code)
            continue
        run = paragraph.add_run(text)
        run.bold = bold
        run.italic = italic
        if code:
            run.font.name = "Consolas"


def _docx_styled_paragraph(doc: Any, style: str, fallback: str) -> Any:
    try:
        return doc.add_paragraph(style=style)
    except KeyError:  # style missing from the template -> nearest base style
        return doc.add_paragraph(style=fallback)


def _docx_code_block(doc: Any, text: str) -> None:
    """One shaded monospace paragraph for a whole code fence (not one per line)."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Pt

    para = _docx_styled_paragraph(doc, "No Spacing", "Normal")
    p_pr = para._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")  # light-grey block fill behind the code
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), "F2F2F2")
    p_pr.append(shd)
    lines = text.split("\n") or [""]
    first = True
    for line in lines:
        run = para.add_run()
        if not first:
            run.add_break()  # keep the block one paragraph, wrap on soft breaks
        run.add_text(line)
        run.font.name = "Consolas"
        run.font.size = Pt(9)
        first = False


def _docx_align(align: str) -> Any:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    return {
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "left": WD_ALIGN_PARAGRAPH.LEFT,
    }.get(align)


def _docx_render(doc: Any, blocks: list[Block]) -> None:
    for b in blocks:
        if b.kind == "heading":
            heading = doc.add_heading(level=b.level)
            _docx_runs(heading, b.runs)  # honor inline runs inside headings
        elif b.kind == "bullet":
            style = "List Bullet" if b.level == 0 else "List Bullet 2"
            _docx_runs(_docx_styled_paragraph(doc, style, "List Bullet"), b.runs)
        elif b.kind == "numbered":
            style = "List Number" if b.level == 0 else "List Number 2"
            _docx_runs(_docx_styled_paragraph(doc, style, "List Number"), b.runs)
        elif b.kind == "code":
            _docx_code_block(doc, b.text)
        elif b.kind == "table":
            cols = max(len(r) for r in b.rows)
            table = doc.add_table(rows=len(b.rows), cols=cols)
            try:
                table.style = "Table Grid"
            except KeyError:
                pass  # borderless is better than no table at all
            for ri, row in enumerate(b.rows):
                for ci in range(cols):
                    cell = table.cell(ri, ci)
                    cell.text = row[ci] if ci < len(row) else ""
                    align = _docx_align(b.aligns[ci]) if ci < len(b.aligns) else None
                    for para in cell.paragraphs:
                        if align is not None:
                            para.alignment = align
                        if ri == 0:  # header row bold
                            for run in para.runs:
                                run.bold = True
        elif b.kind == "hr":
            doc.add_paragraph("─" * 30)
        else:  # paragraph
            _docx_runs(doc.add_paragraph(), b.runs)


# --- xlsx ----------------------------------------------------------------------


def _write_xlsx(
    p: Path,
    content: Any,
    options: "dict[str, Any] | None" = None,
    warnings: "list[str] | None" = None,
) -> None:
    from openpyxl import Workbook

    # BEAUTY is strictly additive here too: without options every statement
    # below is the legacy flow. The one CONTENT extension is chart specs — a
    # sheet value may be {"rows": [...], "charts": [...]} (a shape the legacy
    # writer never accepted, so no existing content changes meaning).
    opts, theme = _beauty_opts(options, warnings)

    wb = Workbook()
    if isinstance(content, dict) and isinstance(content.get("sheets"), dict):
        wb.remove(wb.active)
        for name, rows in content["sheets"].items():
            title = re.sub(r"[\[\]:*?/\\]", "_", str(name))[:31] or "Sheet"
            base, n = title, 1
            while title in wb.sheetnames:  # duplicate names must not crash
                n += 1
                title = f"{base[:28]}~{n}"
            charts = None
            if isinstance(rows, dict) and "rows" in rows:
                charts = rows.get("charts")
                rows = rows.get("rows")
            ws = wb.create_sheet(title=title)
            _xlsx_fill(ws, rows)
            if charts:
                xlsx_add_charts(ws, charts, warnings)
    else:
        ws = wb.active
        if isinstance(content, (list, tuple)):
            for row in content:
                if isinstance(row, (list, tuple)):
                    ws.append([("" if c is None else c) for c in row])
                else:
                    ws.append([row])
        else:
            for line in _as_lines(content):
                ws.append([line])
    if opts:
        for ws in wb.worksheets:
            xlsx_apply_beauty(ws, theme, opts, warnings)
    with _atomic(p) as tmp:
        wb.save(str(tmp))


_XLSX_INT_RX = re.compile(r"-?\d+")
_XLSX_FLOAT_RX = re.compile(r"-?\d+\.\d+")
_XLSX_DATE_RX = re.compile(r"\d{4}-\d{2}-\d{2}")
_XLSX_DATETIME_RX = re.compile(r"\d{4}-\d{2}-\d{2}[ T]\d{2}:\d{2}(:\d{2})?")


def _xlsx_coerce(v: Any) -> tuple[Any, str | None]:
    """Coerce a data cell to number/date; return (value, number_format|None).

    Only clean, unambiguous strings convert: pure ints/floats and ISO dates.
    Formulas (``=``), mixed text, and leading-zero strings (ids, zips, phone
    numbers) stay text so we never silently mangle "007" into 7.
    """
    if not isinstance(v, str):
        return v, None
    s = v.strip()
    if not s or s.startswith("="):
        return v, None
    # A leading zero on a multi-digit integer means "identifier", keep as text.
    if len(s) > 1 and s[0] == "0" and s[1] != ".":
        return v, None
    if _XLSX_INT_RX.fullmatch(s):
        try:
            return int(s), None
        except ValueError:
            return v, None
    if _XLSX_FLOAT_RX.fullmatch(s):
        try:
            return float(s), None
        except ValueError:
            return v, None
    if _XLSX_DATETIME_RX.fullmatch(s):
        try:
            return datetime.fromisoformat(s.replace(" ", "T")), "yyyy-mm-dd hh:mm:ss"
        except ValueError:
            return v, None
    if _XLSX_DATE_RX.fullmatch(s):
        try:
            return datetime.fromisoformat(s), "yyyy-mm-dd"
        except ValueError:
            return v, None
    return v, None


def _xlsx_fill(ws: Any, rows: Any) -> None:
    """Fill one worksheet; bold+freeze a header row, size columns, coerce data."""
    if isinstance(rows, (list, tuple)):
        norm: list[list[Any]] = [
            [("" if c is None else c) for c in row]
            if isinstance(row, (list, tuple))
            else [row]
            for row in rows
        ]
    else:
        norm = [[line] for line in _as_lines(rows)]

    header = (
        len(norm) >= 2
        and bool(norm[0])
        and all(isinstance(c, str) and not c.startswith("=") for c in norm[0])
    )

    for ri, row in enumerate(norm, start=1):
        for ci, val in enumerate(row, start=1):
            cell = ws.cell(row=ri, column=ci)
            if header and ri == 1:
                cell.value = val  # header labels stay verbatim text
                continue
            coerced, fmt = _xlsx_coerce(val)
            cell.value = coerced
            if fmt:
                cell.number_format = fmt

    if header:
        from openpyxl.styles import Font

        for cell in ws[1]:
            cell.font = Font(bold=True)
        ws.freeze_panes = "A2"

    from openpyxl.utils import get_column_letter

    widths: dict[int, int] = {}
    for row in norm:
        for i, c in enumerate(row, start=1):
            widths[i] = max(widths.get(i, 0), len(str(c)))
    for i, w in widths.items():
        ws.column_dimensions[get_column_letter(i)].width = min(max(w + 2, 8), 60)


# --- xlsx beauty (v1.134.0) -----------------------------------------------------
# Public (no underscore) because excel_apply_spec applies the SAME declarative
# options + chart specs — one implementation, one validation, one behavior.

#: A1-style range/cell — validated up front so a model's typo becomes a
#: recorded warning, never an openpyxl traceback mid-write.
_A1_RANGE_RX = re.compile(r"^\$?[A-Za-z]{1,3}\$?\d{1,7}(:\$?[A-Za-z]{1,3}\$?\d{1,7})?$")
_A1_CELL_RX = re.compile(r"^\$?[A-Za-z]{1,3}\$?\d{1,7}$")
_COL_LETTER_RX = re.compile(r"^[A-Za-z]{1,3}$")

_CHART_TYPES = ("bar", "line", "pie")


def _chart_bounds(ws: Any, ref: Any) -> "tuple[int, int, int, int] | None":
    """(min_col, min_row, max_col, max_row) for an A1 ref that is well-formed
    AND inside the sheet's used dimensions; None otherwise."""
    from openpyxl.utils import range_boundaries

    s = str(ref or "").strip()
    if not _A1_RANGE_RX.fullmatch(s):
        return None
    try:
        min_col, min_row, max_col, max_row = range_boundaries(s.upper())
    except Exception:  # noqa: BLE001 — malformed refs must degrade, not crash
        return None
    if None in (min_col, min_row, max_col, max_row):
        return None
    if min_col > max_col or min_row > max_row:
        return None
    if max_row > ws.max_row or max_col > ws.max_column:
        return None
    return min_col, min_row, max_col, max_row


def xlsx_add_charts(ws: Any, charts: Any, warnings: "list[str] | None") -> int:
    """Add declared charts to ``ws``; returns how many landed.

    Each spec: {"type": "bar"|"line"|"pie", "title": str, "data_range":
    "B2:B10", "categories_range": "A2:A10", "anchor": "E2"}. Anything invalid
    (unknown type, malformed/out-of-bounds range, bad anchor) SKIPS that chart
    with a recorded warning — a chart typo must never cost the workbook.
    """
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

    if not isinstance(charts, (list, tuple)):
        _warn(warnings, f"charts must be a list of chart specs on sheet {ws.title!r}")
        return 0
    added = 0
    for i, spec in enumerate(charts, start=1):
        label = f"chart {i} on sheet {ws.title!r}"
        if not isinstance(spec, dict):
            _warn(warnings, f"{label} skipped: spec must be an object")
            continue
        ctype = str(spec.get("type") or "").strip().lower()
        if ctype not in _CHART_TYPES:
            _warn(
                warnings,
                f"{label} skipped: unknown type {spec.get('type')!r} "
                f"(use one of {', '.join(_CHART_TYPES)})",
            )
            continue
        data = _chart_bounds(ws, spec.get("data_range"))
        if data is None:
            _warn(
                warnings,
                f"{label} skipped: data_range {spec.get('data_range')!r} is not "
                f"an A1 range inside the used cells "
                f"(sheet spans A1:{_col_letter(ws.max_column)}{ws.max_row})",
            )
            continue
        cats = None
        if str(spec.get("categories_range") or "").strip():
            cats = _chart_bounds(ws, spec.get("categories_range"))
            if cats is None:
                _warn(
                    warnings,
                    f"{label} skipped: categories_range "
                    f"{spec.get('categories_range')!r} is not an A1 range "
                    f"inside the used cells",
                )
                continue
        anchor = str(spec.get("anchor") or "").strip().upper() or "E2"
        if not _A1_CELL_RX.fullmatch(anchor):
            _warn(warnings, f"{label} skipped: anchor {spec.get('anchor')!r} is not a cell like 'E2'")
            continue
        chart = {"bar": BarChart, "line": LineChart, "pie": PieChart}[ctype]()
        title = str(spec.get("title") or "").strip()
        if title:
            chart.title = title
        chart.add_data(
            Reference(ws, min_col=data[0], min_row=data[1], max_col=data[2], max_row=data[3]),
            titles_from_data=False,
        )
        if cats is not None:
            chart.set_categories(
                Reference(ws, min_col=cats[0], min_row=cats[1], max_col=cats[2], max_row=cats[3])
            )
        ws.add_chart(chart, anchor.replace("$", ""))
        added += 1
    return added


def _col_letter(idx: int) -> str:
    from openpyxl.utils import get_column_letter

    return get_column_letter(max(int(idx), 1))


def _xlsx_header_like(ws: Any) -> bool:
    """Same heuristic as ``_xlsx_fill``: >=2 rows and a first row that is all
    non-formula text reads as a header row."""
    if ws.max_row < 2:
        return False
    first = [c.value for c in ws[1]]
    return bool(first) and all(
        isinstance(v, str) and not v.startswith("=") for v in first
    )


def xlsx_apply_beauty(
    ws: Any,
    theme: "Theme | None",
    opts: dict[str, Any],
    warnings: "list[str] | None",
) -> None:
    """Apply the declarative beauty options on top of a filled worksheet.

    * ``autosize`` (default ON when a theme is set): width from the longest
      cell text per column, clamped 8..60.
    * ``freeze_header``: freeze row 1.
    * themed header row (fill + contrast font) when a header row is detected.
    * ``banded``: alternate data rows get the theme band fill (or a neutral
      grey without a theme).
    * ``number_formats``: {column letter: Excel format} on data rows.
    """
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import column_index_from_string

    max_row, max_col = ws.max_row, ws.max_column
    header = _xlsx_header_like(ws)

    if bool(opts.get("autosize", theme is not None)):
        widths: dict[int, int] = {}
        for row in ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col):
            for c in row:
                if c.value is not None:
                    widths[c.column] = max(widths.get(c.column, 0), len(str(c.value)))
        for ci, w in widths.items():
            ws.column_dimensions[_col_letter(ci)].width = min(max(w + 2, 8), 60)

    if bool(opts.get("freeze_header")):
        ws.freeze_panes = "A2"

    if theme is not None and header:
        head_fill = PatternFill(
            start_color=theme.table_header_fill,
            end_color=theme.table_header_fill,
            fill_type="solid",
        )
        head_font = Font(bold=True, color=theme.table_header_font)
        for c in ws[1]:
            c.fill = head_fill
            c.font = head_font

    if bool(opts.get("banded")) and max_row >= 2:
        band_rgb = theme.band_fill if theme is not None else "F2F2F2"
        band = PatternFill(start_color=band_rgb, end_color=band_rgb, fill_type="solid")
        # Shade every second data row (3, 5, ...) so the row right under the
        # header stays clean — the classic banded-table look.
        for ri in range(3, max_row + 1, 2):
            for ci in range(1, max_col + 1):
                ws.cell(row=ri, column=ci).fill = band

    fmts = opts.get("number_formats")
    if isinstance(fmts, dict):
        for letter, fmt in fmts.items():
            col = str(letter or "").strip().upper()
            if not _COL_LETTER_RX.fullmatch(col):
                _warn(
                    warnings,
                    f"number_formats key {letter!r} skipped on sheet "
                    f"{ws.title!r}: use a column letter like 'B'",
                )
                continue
            ci = column_index_from_string(col)
            if ci > max_col:
                continue  # a format for an unused column is a harmless no-op
            first = 2 if header else 1  # header labels keep General
            for ri in range(first, max_row + 1):
                ws.cell(row=ri, column=ci).number_format = str(fmt)
    elif fmts is not None:
        _warn(warnings, "number_formats must be {column letter: format string} — ignored")


# --- pptx ----------------------------------------------------------------------

#: Bullets past this many spill onto a "Title (cont.)" continuation slide so a
#: dense section never runs silently off the bottom of the slide.
_PPTX_MAX_ITEMS = 9


def _write_pptx(p: Path, content: Any) -> None:
    import pptx

    sections = (
        _pptx_sections(parse_markdown(content)) if isinstance(content, str) else None
    )

    prs = pptx.Presentation()

    if sections is None:
        # Legacy flat deck: title slide + one bullet slide with every line.
        lines = list(_as_lines(content))
        title = lines[0] if lines and lines[0].strip() else "Document"

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Generated by Iron Jarvis"

        bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
        bullet_slide.shapes.title.text = title
        body = bullet_slide.placeholders[1].text_frame
        _pptx_prep_body(body)
        first, *rest = lines or [""]
        body.text = first
        for line in rest:
            body.add_paragraph().text = line
    else:
        # Sectioned deck: title slide + one (or more) slide per '# '/'---'.
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = sections[0][0] or "Document"
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Generated by Iron Jarvis"

        for title, items, tables, notes in sections:
            _pptx_content_slides(prs, title or "Section", items, tables, notes)

    with _atomic(p) as tmp:
        prs.save(str(tmp))


def _pptx_prep_body(tf: Any) -> None:
    """Make a body text frame wrap + shrink-to-fit instead of overflowing."""
    from pptx.enum.text import MSO_AUTO_SIZE

    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:  # pragma: no cover - some templates reject it
        pass


def _pptx_fill_para(para: Any, runs: list[Run]) -> None:
    """Render styled runs onto one bullet paragraph (bold/italic/code/link)."""
    for r in runs:
        text, bold, italic, code, href, _img = _run_parts(r)
        run = para.add_run()
        run.text = text
        if bold:
            run.font.bold = True
        if italic:
            run.font.italic = True
        if code:
            run.font.name = "Consolas"
        if href:
            try:
                run.hyperlink.address = href
            except Exception:  # pragma: no cover - defensive
                pass


def _pptx_content_slides(
    prs: Any,
    title: str,
    items: list[tuple[list[Run], int]],
    tables: list[Block],
    notes: str,
) -> None:
    """Emit bullet slide(s) (+ continuation slides) then any table slides."""
    chunks = [
        items[i : i + _PPTX_MAX_ITEMS]
        for i in range(0, max(len(items), 1), _PPTX_MAX_ITEMS)
    ]
    first_slide = None
    for ci, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title if ci == 0 else f"{title} (cont.)"
        first_slide = first_slide or slide
        body = slide.placeholders[1].text_frame
        _pptx_prep_body(body)
        for idx, (runs, level) in enumerate(chunk or [([Run("", False, False)], 0)]):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.level = min(level, 4)
            _pptx_fill_para(para, runs)

    for tb in tables:
        _pptx_table_slide(prs, f"{title} (cont.)", tb)

    if notes and first_slide is not None:
        try:
            first_slide.notes_slide.notes_text_frame.text = notes
        except Exception:  # pragma: no cover - defensive
            pass


def _pptx_table_slide(prs: Any, title: str, b: Block) -> None:
    """Render a table Block as a REAL pptx table on its own slide."""
    from pptx.util import Emu, Inches

    try:
        layout = prs.slide_layouts[5]  # "Title Only" in the default template
    except IndexError:  # pragma: no cover
        layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = title
    except AttributeError:  # pragma: no cover
        pass

    rows = b.rows
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    left, top = Inches(0.5), Inches(1.5)
    width = prs.slide_width - Emu(2 * Inches(0.5))
    height = Inches(0.4) * n_rows
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            tbl.cell(ri, ci).text = row[ci] if ci < len(row) else ""


def _pptx_sections(
    blocks: list[Block],
) -> list[tuple[str, list[tuple[list[Run], int]], list[Block], str]] | None:
    """Split blocks into ``(title, items, tables, notes)`` sections.

    A level-1 heading or a ``---`` rule starts a new section. ``items`` are
    ``(runs, indent)`` bullet lines; tables are kept as Blocks for real rendering;
    a ``Notes:`` paragraph feeds the slide's speaker notes. Returns ``None`` when
    the content has neither heading nor rule, so the caller keeps the flat deck.
    """
    if not any(
        (b.kind == "heading" and b.level == 1) or b.kind == "hr" for b in blocks
    ):
        return None

    sections: list[tuple[str, list[tuple[list[Run], int]], list[Block], list[str]]] = []

    def cur() -> tuple[str, list, list, list]:
        if not sections:
            sections.append(("", [], [], []))
        return sections[-1]

    for b in blocks:
        if b.kind == "heading" and b.level == 1:
            # An hr immediately followed by a heading is ONE section, titled.
            if sections and sections[-1][0] == "" and not any(sections[-1][1:]):
                sections[-1] = (b.text, sections[-1][1], sections[-1][2], sections[-1][3])
            else:
                sections.append((b.text, [], [], []))
            continue
        if b.kind == "hr":
            sections.append(("", [], [], []))
            continue

        _title, items, tables, notes = cur()
        if b.kind in ("bullet", "numbered"):
            items.append((b.runs, b.level))
        elif b.kind == "table":
            tables.append(b)
        elif b.kind == "code":
            for code_line in b.text.split("\n"):
                items.append(([Run(code_line, code=True)], 1))
        elif b.kind == "paragraph" and b.text.startswith("Notes:"):
            notes.append(b.text[len("Notes:") :].strip())
        else:  # paragraph / sub-heading
            items.append((b.runs, 0))

    # A trailing '---' (footer rule) must not yield a blank slide.
    while sections and sections[-1][0] == "" and not any(sections[-1][1:]):
        sections.pop()
    if not sections:
        return None
    return [(t, it, tb, "\n".join(nt)) for (t, it, tb, nt) in sections]


# --- pdf -----------------------------------------------------------------------


_PDF_HEADING_SIZES = {1: 20, 2: 16, 3: 14, 4: 12}


def _latin1(text: str) -> str:
    # Core fonts are Latin-1 only; replace anything outside it so fpdf2 cannot crash.
    return text.encode("latin-1", "replace").decode("latin-1")


#: Bundled DejaVu TTFs (Bitstream Vera license) for FULL-UNICODE PDF output —
#: accents, Cyrillic, Greek, symbols. When present, PDFs keep real unicode;
#: when missing (stripped install), we fall back to the core Latin-1 fonts with
#: the historical 'replace' sanitiser so nothing ever crashes.
_FONT_DIR = Path(__file__).resolve().parent / "fonts"


def _pdf_fonts(pdf: Any) -> tuple[str, str, Any]:
    """Register unicode fonts on ``pdf``; return (sans, mono, sanitize)."""
    try:
        sans = _FONT_DIR / "DejaVuSans.ttf"
        bold = _FONT_DIR / "DejaVuSans-Bold.ttf"
        mono = _FONT_DIR / "DejaVuSansMono.ttf"
        if sans.is_file() and bold.is_file() and mono.is_file():
            pdf.add_font("DJSans", "", str(sans))
            pdf.add_font("DJSans", "B", str(bold))
            pdf.add_font("DJMono", "", str(mono))
            return "DJSans", "DJMono", lambda s: s  # true unicode — no mangling
    except Exception:  # noqa: BLE001 — font trouble must never break a write
        pass
    return "Helvetica", "Courier", _latin1


def _pdf_usable(pdf: Any) -> float:
    return pdf.w - pdf.l_margin - pdf.r_margin


def _prebreak(pdf: Any, text: str, max_w: float) -> str:
    """Hard-break any single token wider than ``max_w`` so wrapping can't fail.

    fpdf2's ``multi_cell`` raises when one unbreakable word is wider than the
    cell; we slice such words across newlines (which multi_cell honours) up
    front. Requires the current font to already be set (widths depend on it).
    """
    if max_w <= 0:
        return text
    out: list[str] = []
    for word in text.split(" "):
        if not word or pdf.get_string_width(word) <= max_w:
            out.append(word)
            continue
        chunks: list[str] = []
        cur = ""
        for ch in word:
            if cur and pdf.get_string_width(cur + ch) > max_w:
                chunks.append(cur)
                cur = ch
            else:
                cur += ch
        if cur:
            chunks.append(cur)
        out.append("\n".join(chunks))
    return " ".join(out)


def _write_pdf(p: Path, content: Any) -> None:
    """Compose a PDF in memory, then write atomically.

    Building in memory (``pdf.output()`` -> bytes) means a layout failure never
    leaves a truncated file. If a render still trips an ``FPDFException`` we
    retry smaller, then as hard-chunked plain text, so SOME valid file always
    lands rather than none.
    """
    from fpdf.errors import FPDFException

    data: bytes | None = None
    for scale, plain in ((1.0, False), (0.7, False), (0.5, True)):
        try:
            pdf = _compose_pdf(content, scale, plain)
            data = bytes(pdf.output())
            break
        except FPDFException:
            continue
    if data is None:  # last resort: a minimal but valid one-page document
        pdf = _compose_pdf(" ", 0.5, True)
        data = bytes(pdf.output())
    with _atomic(p) as tmp:
        Path(tmp).write_bytes(data)


def _compose_pdf(content: Any, scale: float, plain: bool) -> Any:
    from fpdf import FPDF

    pdf = FPDF()
    sans, mono, clean = _pdf_fonts(pdf)
    pdf.add_page()
    if plain or not isinstance(content, str):
        pdf.set_font(sans, size=max(6, int(12 * scale)))
        text = clean(_as_text(content))
        if not text.strip():
            text = " "
        pdf.multi_cell(0, 8 * scale, _prebreak(pdf, text, _pdf_usable(pdf) - 1))
        return pdf
    _pdf_render(pdf, parse_markdown(content), sans, mono, clean, scale)
    return pdf


def _pdf_render(
    pdf: Any, blocks: list[Block], sans: str, mono: str, clean: Any, scale: float = 1.0
) -> None:
    if not blocks:  # keep the empty-content page valid, as before
        pdf.set_font(sans, size=12)
        pdf.multi_cell(0, 8, " ")
        return
    number = 0  # running counter for consecutive numbered items
    for b in blocks:
        if b.kind != "numbered":
            number = 0
        usable = _pdf_usable(pdf)
        if b.kind == "heading":
            size = max(7, int(_PDF_HEADING_SIZES.get(b.level, 12) * scale))
            pdf.set_font(sans, "B", size)
            pdf.multi_cell(0, size * 0.5 + 2, _prebreak(pdf, clean(b.text) or " ", usable - 1))
            pdf.ln(2)
        elif b.kind == "bullet":
            pdf.set_font(sans, size=max(7, int(11 * scale)))
            indent = 5 * (b.level + 1)
            pdf.set_x(pdf.l_margin + indent)
            pdf.multi_cell(0, 6, _prebreak(pdf, "- " + clean(b.text), usable - indent - 1))
        elif b.kind == "numbered":
            number += 1
            pdf.set_font(sans, size=max(7, int(11 * scale)))
            indent = 5 * (b.level + 1)
            pdf.set_x(pdf.l_margin + indent)
            pdf.multi_cell(
                0, 6, _prebreak(pdf, f"{number}. " + clean(b.text), usable - indent - 1)
            )
        elif b.kind == "code":
            pdf.set_font(mono, size=max(6, int(10 * scale)))
            pdf.set_fill_color(235, 235, 235)
            for line in b.text.split("\n"):
                pdf.multi_cell(0, 5, _prebreak(pdf, clean(line) or " ", usable - 1), fill=True)
            pdf.ln(2)
        elif b.kind == "table":
            _pdf_table(pdf, b.rows, b.aligns, sans, clean, scale)
        elif b.kind == "hr":
            y = pdf.get_y() + 2
            pdf.line(pdf.l_margin, y, pdf.w - pdf.r_margin, y)
            pdf.set_y(y + 3)
        else:  # paragraph
            pdf.set_font(sans, size=max(7, int(11 * scale)))
            pdf.multi_cell(0, 6, _prebreak(pdf, clean(b.text) or " ", usable - 1))
            pdf.ln(1)


_PDF_ALIGN = {"left": "LEFT", "right": "RIGHT", "center": "CENTER"}


def _pdf_table(
    pdf: Any, rows: list[list[str]], aligns: list[str], sans: str, clean: Any, scale: float
) -> None:
    cols = max(len(r) for r in rows)
    pdf.set_font(sans, size=max(7, int(10 * scale)))
    if not hasattr(pdf, "table"):  # older fpdf2 -> hand-wrapped fallback
        _pdf_table_fallback(pdf, rows, cols, sans, clean, scale)
        pdf.ln(2)
        return
    text_align = tuple(
        _PDF_ALIGN.get(aligns[i] if i < len(aligns) else "", "LEFT") for i in range(cols)
    )
    try:
        # fpdf2's native table WRAPS long cell text, sizes columns, and repeats
        # the heading row after a page break — none of which the old fixed-width
        # ``cell`` grid did (it clipped and mis-aligned on overflow).
        with pdf.table(text_align=text_align, first_row_as_headings=True) as table:
            for row in rows:
                trow = table.row()
                for ci in range(cols):
                    trow.cell(clean(row[ci]) if ci < len(row) else "")
    except Exception:  # noqa: BLE001 — any table quirk degrades to the fallback
        _pdf_table_fallback(pdf, rows, cols, sans, clean, scale)
    pdf.ln(2)


def _pdf_table_fallback(
    pdf: Any, rows: list[list[str]], cols: int, sans: str, clean: Any, scale: float
) -> None:
    """Wrapping grid for fpdf2 builds without ``pdf.table`` (multi_cell per cell)."""
    width = _pdf_usable(pdf) / max(cols, 1)
    for ri, row in enumerate(rows):
        pdf.set_font(sans, "B" if ri == 0 else "", max(7, int(10 * scale)))
        x0, y0 = pdf.get_x(), pdf.get_y()
        max_y = y0
        for ci in range(cols):
            cell = clean(row[ci]) if ci < len(row) else ""
            pdf.set_xy(x0 + ci * width, y0)
            pdf.multi_cell(width, 6, _prebreak(pdf, cell, width - 1), border=1)
            max_y = max(max_y, pdf.get_y())
        pdf.set_xy(x0, max_y)


# --- csv / text ------------------------------------------------------------------


def _write_csv(p: Path, content: Any) -> None:
    with _atomic(p) as tmp:
        with open(tmp, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if isinstance(content, (list, tuple)):
                for row in content:
                    if isinstance(row, (list, tuple)):
                        writer.writerow(list(row))
                    else:
                        writer.writerow([row])
            else:
                for line in _as_lines(content):
                    writer.writerow([line])


def _write_text(p: Path, content: Any) -> None:
    _atomic_text(p, _as_text(content))


# --- html ----------------------------------------------------------------------


_HTML_HEAD = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
body { font-family: -apple-system, "Segoe UI", Arial, sans-serif; color: #1a1a1a;
       background: #ffffff; max-width: 800px; margin: 2rem auto; padding: 0 1rem;
       line-height: 1.5; }
h1, h2, h3, h4 { line-height: 1.25; }
pre { background: #f4f4f4; padding: 0.75rem; overflow-x: auto; }
code, pre { font-family: Consolas, "Courier New", monospace; font-size: 0.9em; }
table { border-collapse: collapse; margin: 1em 0; }
th, td { border: 1px solid #999; padding: 0.35em 0.6em; text-align: left; }
th { background: #f0f0f0; }
img { max-width: 100%; }
hr { border: none; border-top: 1px solid #ccc; margin: 1.5em 0; }
</style>
</head>
<body>
"""

_HTML_FOOT = "\n</body>\n</html>\n"


def html_page(markdown_text: str, title: str = "") -> str:
    """A full standalone HTML page rendered from markdown TEXT (nothing is
    written to disk) — the same rendering ``write_document`` uses for ``.html``
    files, so shared/exported pages look identical to created ones."""
    body = _html_render(parse_markdown(markdown_text))
    head = _HTML_HEAD
    if title:
        head = head.replace(
            "<style>", f"<title>{_html_mod.escape(title)}</title>\n<style>", 1
        )
    return head + body + _HTML_FOOT


def _write_html(p: Path, content: Any) -> None:
    if isinstance(content, str):
        sniff = content.lstrip().lower()
        if sniff.startswith("<!doctype") or sniff.startswith("<html"):
            _atomic_text(p, content)  # already a full page
            return
        body = _html_render(parse_markdown(content))
    elif isinstance(content, (list, tuple)) and any(
        isinstance(r, (list, tuple)) for r in content
    ):
        rows = [
            [str(c) for c in r] if isinstance(r, (list, tuple)) else [str(r)]
            for r in content
        ]
        body = _html_table(rows, [])
    else:
        body = _html_render(parse_markdown(_as_text(content)))
    _atomic_text(p, _HTML_HEAD + body + _HTML_FOOT)


def _runs_html(runs: list[Run]) -> str:
    parts: list[str] = []
    for r in runs:
        text, bold, italic, code, href, image = _run_parts(r)
        if image and href:
            src = _html_mod.escape(href, quote=True)
            alt = _html_mod.escape(text)
            parts.append(f'<img src="{src}" alt="{alt}">')
            continue
        chunk = _html_mod.escape(text)
        if code:
            chunk = f"<code>{chunk}</code>"
        if bold:
            chunk = f"<strong>{chunk}</strong>"
        if italic:
            chunk = f"<em>{chunk}</em>"
        if href:
            chunk = f'<a href="{_html_mod.escape(href, quote=True)}">{chunk}</a>'
        parts.append(chunk)
    return "".join(parts)


def _html_align_attr(align: str) -> str:
    return f' style="text-align:{align}"' if align in ("center", "right", "left") else ""


def _html_table(rows: list[list[str]], aligns: list[str]) -> str:
    parts = ["<table>"]
    for ri, row in enumerate(rows):
        tag = "th" if ri == 0 else "td"
        cells = "".join(
            f"<{tag}{_html_align_attr(aligns[ci] if ci < len(aligns) else '')}>"
            f"{_html_mod.escape(c)}</{tag}>"
            for ci, c in enumerate(row)
        )
        parts.append(f"<tr>{cells}</tr>")
    parts.append("</table>")
    return "\n".join(parts)


def _html_lists(items: list[Block]) -> str:
    """Render a run of list Blocks with VALID nesting.

    A deeper item's sublist is wrapped in its own ``<li>`` (``<li><ul>…</ul></li>``)
    rather than the old invalid ``<ul><ul>``; ``ul`` and ``ol`` open/close
    independently by tracking each item's kind at each level.
    """
    pos = 0

    def render(level: int) -> str:
        nonlocal pos
        out: list[str] = []
        while pos < len(items) and items[pos].level >= level:
            kind = items[pos].kind
            tag = "ul" if kind == "bullet" else "ol"
            out.append(f"<{tag}>")
            while (
                pos < len(items)
                and items[pos].level >= level
                and items[pos].kind == kind
            ):
                if items[pos].level > level:
                    # A deeper run is a child list — wrap it in its OWN <li>
                    # (valid ``<li><ul>…</ul></li>`` instead of the old
                    # invalid ``<ul><ul>``), independent of ul/ol kind.
                    out.append(f"<li>{render(items[pos].level)}</li>")
                    continue
                out.append(f"<li>{_runs_html(items[pos].runs)}</li>")
                pos += 1
            out.append(f"</{tag}>")
        return "".join(out)

    return render(items[0].level)


def _html_render(blocks: list[Block]) -> str:
    out: list[str] = []
    i = 0
    while i < len(blocks):
        b = blocks[i]
        if b.kind in ("bullet", "numbered"):
            items: list[Block] = []
            while i < len(blocks) and blocks[i].kind in ("bullet", "numbered"):
                items.append(blocks[i])
                i += 1
            out.append(_html_lists(items))
            continue
        if b.kind == "heading":
            out.append(f"<h{b.level}>{_runs_html(b.runs)}</h{b.level}>")
        elif b.kind == "code":
            out.append(f"<pre><code>{_html_mod.escape(b.text)}</code></pre>")
        elif b.kind == "table":
            out.append(_html_table(b.rows, b.aligns))
        elif b.kind == "hr":
            out.append("<hr>")
        else:  # paragraph
            out.append(f"<p>{_runs_html(b.runs)}</p>")
        i += 1
    return "\n".join(out)
